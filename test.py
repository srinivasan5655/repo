import streamlit as st
import zipfile
import os
import tempfile
from pathlib import Path
from openai import AzureOpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
import json
from typing import Dict, List, Tuple
import re
from datetime import datetime

# Page config
st.set_page_config(
    page_title="Business Proposal & Architecture Generator",
    page_icon="📊",
    layout="wide"
)

# Initialize session state
if 'analysis_complete' not in st.session_state:
    st.session_state.analysis_complete = False
if 'analysis_data' not in st.session_state:
    st.session_state.analysis_data = None

class CodebaseAnalyzer:
    """Analyzes codebase and generates architecture insights"""
    
    def __init__(self, azure_endpoint: str, azure_key: str, azure_deployment: str):
        self.client = AzureOpenAI(
            azure_endpoint=azure_endpoint,
            api_key=azure_key,
            api_version="2024-05-01-preview"
        )
        self.deployment = azure_deployment
        
    def extract_files(self, zip_path: str, extract_to: str) -> Dict[str, str]:
        """Extract and read files from zip"""
        file_contents = {}
        
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(extract_to)
            
        # Read all code files
        for root, dirs, files in os.walk(extract_to):
            dirs[:] = [d for d in dirs if d not in ['.git', 'node_modules', '__pycache__', 'venv', '.venv', 'dist', 'build']]
            
            for file in files:
                if self._is_code_file(file):
                    file_path = os.path.join(root, file)
                    try:
                        with open(file_path, 'r', encoding='utf-8') as f:
                            relative_path = os.path.relpath(file_path, extract_to)
                            file_contents[relative_path] = f.read()
                    except Exception as e:
                        st.warning(f"Could not read {file}: {e}")
                        
        return file_contents
    
    def _is_code_file(self, filename: str) -> bool:
        """Check if file is a code file"""
        code_extensions = [
            '.py', '.js', '.jsx', '.ts', '.tsx', '.java', '.cpp', '.c', '.h',
            '.cs', '.go', '.rs', '.rb', '.php', '.swift', '.kt', '.scala',
            '.json', '.yaml', '.yml', '.xml', '.sql', '.sh', '.md', '.txt'
        ]
        return any(filename.endswith(ext) for ext in code_extensions)
    
    def analyze_codebase(self, file_contents: Dict[str, str], progress_callback=None) -> Dict:
        """Analyze codebase using Azure OpenAI"""
        
        file_summary = "\n".join([f"File: {path}\nLines: {len(content.splitlines())}" 
                                  for path, content in list(file_contents.items())[:50]])
        
        limited_contents = {}
        total_chars = 0
        max_chars = 50000
        
        for path, content in file_contents.items():
            if total_chars + len(content) > max_chars:
                break
            limited_contents[path] = content[:5000]
            total_chars += len(content)
        
        analyses = {}
        
        # 1. Business Value & Executive Summary
        if progress_callback:
            progress_callback("Analyzing business value...")
        
        business_prompt = f"""Analyze this codebase from a BUSINESS PERSPECTIVE and provide:
1. Main business problem it solves
2. Key value propositions
3. Target users/stakeholders
4. Business benefits and ROI potential
5. Competitive advantages

Codebase structure:
{file_summary}

Respond in JSON format:
{{
    "business_problem": "...",
    "value_propositions": ["..."],
    "target_users": ["..."],
    "business_benefits": ["..."],
    "competitive_advantages": ["..."],
    "executive_summary": "..."
}}"""
        
        analyses['business'] = self._call_llm(business_prompt)
        
        # 2. Overall Architecture
        if progress_callback:
            progress_callback("Analyzing technical architecture...")
        
        arch_prompt = f"""Analyze this codebase structure and provide:
1. Main architecture pattern (MVC, Microservices, Layered, etc.)
2. Key components and their responsibilities
3. Technology stack
4. Main modules/packages
5. Scalability characteristics

Codebase structure:
{file_summary}

Sample files:
{self._format_files_for_prompt(limited_contents, 3)}

Respond in JSON format:
{{
    "architecture_pattern": "...",
    "technology_stack": ["..."],
    "main_components": [
        {{"name": "...", "responsibility": "...", "type": "frontend/backend/database/service"}}
    ],
    "modules": [
        {{"name": "...", "purpose": "..."}}
    ],
    "scalability": "description of scalability"
}}"""
        
        analyses['architecture'] = self._call_llm(arch_prompt)
        
        # 3. Technical Capabilities
        if progress_callback:
            progress_callback("Analyzing technical capabilities...")
            
        capabilities_prompt = f"""Analyze the technical capabilities and features:

Files:
{self._format_files_for_prompt(limited_contents, 5)}

Provide in JSON:
{{
    "core_features": ["..."],
    "technical_capabilities": ["..."],
    "performance_characteristics": ["..."],
    "integration_points": ["..."]
}}"""
        
        analyses['capabilities'] = self._call_llm(capabilities_prompt)
        
        # 4. Implementation Roadmap
        if progress_callback:
            progress_callback("Creating implementation roadmap...")
            
        roadmap_prompt = f"""Based on this codebase, suggest an implementation/deployment roadmap:

Files:
{self._format_files_for_prompt(limited_contents, 3)}

Provide in JSON:
{{
    "phases": [
        {{"phase": "Phase 1", "title": "...", "duration": "...", "deliverables": ["..."]}}
    ],
    "milestones": ["..."],
    "success_criteria": ["..."]
}}"""
        
        analyses['roadmap'] = self._call_llm(roadmap_prompt)
        
        # 5. Dependencies & Integrations
        if progress_callback:
            progress_callback("Analyzing dependencies...")
            
        dep_prompt = f"""Analyze the dependencies and integrations:

Files:
{self._format_files_for_prompt(limited_contents, 5)}

Provide in JSON:
{{
    "dependencies": [
        {{"from": "ComponentA", "to": "ComponentB", "type": "uses/imports/calls"}}
    ],
    "external_integrations": ["..."],
    "third_party_services": ["..."]
}}"""
        
        analyses['dependencies'] = self._call_llm(dep_prompt)
        
        # 6. Security & Compliance
        if progress_callback:
            progress_callback("Analyzing security measures...")
            
        security_prompt = f"""Analyze security and compliance aspects:

Files:
{self._format_files_for_prompt(limited_contents, 3)}

Provide in JSON:
{{
    "security_measures": ["..."],
    "compliance_standards": ["..."],
    "data_protection": ["..."],
    "authentication_methods": ["..."]
}}"""
        
        analyses['security'] = self._call_llm(security_prompt)
        
        # 7. Cost & Resource Estimation
        if progress_callback:
            progress_callback("Estimating resources...")
            
        cost_prompt = f"""Estimate infrastructure and resource requirements:

Technology stack: {analyses.get('architecture', {}).get('technology_stack', [])}

Provide in JSON:
{{
    "infrastructure_needs": ["..."],
    "team_requirements": ["..."],
    "estimated_timeline": "...",
    "maintenance_considerations": ["..."]
}}"""
        
        analyses['resources'] = self._call_llm(cost_prompt)
        
        return analyses
    
    def _format_files_for_prompt(self, files: Dict[str, str], limit: int) -> str:
        """Format files for LLM prompt"""
        result = []
        for i, (path, content) in enumerate(list(files.items())[:limit]):
            result.append(f"\n--- {path} ---\n{content[:2000]}\n")
        return "\n".join(result)
    
    def _call_llm(self, prompt: str) -> Dict:
        """Call Azure OpenAI and parse JSON response"""
        try:
            response = self.client.chat.completions.create(
                model=self.deployment,
                messages=[
                    {"role": "system", "content": "You are a business and technical analyst. Always respond with valid JSON."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                max_tokens=2000
            )
            
            content = response.choices[0].message.content
            
            json_match = re.search(r'```json\s*(.*?)\s*```', content, re.DOTALL)
            if json_match:
                content = json_match.group(1)
            
            return json.loads(content)
        except Exception as e:
            st.error(f"LLM call failed: {e}")
            return {}


class BusinessPPTGenerator:
    """Generates stunning business proposal presentations"""
    
    # Professional color palette
    COLORS = {
        'primary': RGBColor(20, 33, 61),        # Deep Navy
        'secondary': RGBColor(0, 123, 255),     # Vibrant Blue
        'accent': RGBColor(255, 107, 53),       # Coral Orange
        'success': RGBColor(40, 167, 69),       # Professional Green
        'gold': RGBColor(255, 193, 7),          # Gold
        'purple': RGBColor(111, 66, 193),       # Royal Purple
        'light': RGBColor(248, 249, 250),       # Off White
        'text': RGBColor(33, 37, 41),           # Dark Gray
        'text_light': RGBColor(108, 117, 125),  # Medium Gray
        'border': RGBColor(222, 226, 230),      # Light Border
        'white': RGBColor(255, 255, 255),
    }
    
    def __init__(self, company_name: str = "Your Company"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        self.prs.slide_height = Inches(7.5)
        self.company_name = company_name
        
    def _add_premium_background(self, slide, style='default'):
        """Add premium background with modern design"""
        # Base background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['white']
        
        if style == 'cover':
            # Diagonal split background
            shape1 = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(7), Inches(7.5)
            )
            fill = shape1.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['primary']
            shape1.line.fill.background()
            
            shape2 = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(7), Inches(0), Inches(6.333), Inches(7.5)
            )
            fill = shape2.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['secondary']
            shape2.line.fill.background()
            
        elif style == 'section':
            # Gradient-style background with shapes
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(13.333), Inches(7.5)
            )
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['primary']
            bg.line.fill.background()
            
            # Decorative circles
            for i, (x, y, size, color) in enumerate([
                (10, -1, 4, self.COLORS['secondary']),
                (-1, 5, 3.5, self.COLORS['accent']),
                (11, 6, 2.5, self.COLORS['gold'])
            ]):
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(x), Inches(y), Inches(size), Inches(size)
                )
                fill = circle.fill
                fill.solid()
                fill.fore_color.rgb = color
                circle.fill.transparency = 0.3
                circle.line.fill.background()
        
        else:  # content slides
            # Subtle accent bar on left
            accent_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(0.15), Inches(7.5)
            )
            fill = accent_bar.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['secondary']
            accent_bar.line.fill.background()
            
            # Top decoration
            top_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(11.5), Inches(-0.5), Inches(2.5), Inches(2)
            )
            fill = top_shape.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['light']
            top_shape.line.fill.background()
    
    def add_cover_slide(self, title: str, subtitle: str):
        """Add stunning cover slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide, 'cover')
        
        # Title on dark side
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(5.5), Inches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        title_para.line_spacing = 1.1
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(5.5), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(22)
        subtitle_para.font.color.rgb = self.COLORS['white']
        subtitle_para.font.italic = True
        
        # Company name and date on light side
        company_box = slide.shapes.add_textbox(Inches(7.5), Inches(3), Inches(5), Inches(0.6))
        company_frame = company_box.text_frame
        company_frame.text = self.company_name
        company_para = company_frame.paragraphs[0]
        company_para.font.size = Pt(28)
        company_para.font.bold = True
        company_para.font.color.rgb = self.COLORS['white']
        
        date_box = slide.shapes.add_textbox(Inches(7.5), Inches(3.7), Inches(5), Inches(0.4))
        date_frame = date_box.text_frame
        date_frame.text = datetime.now().strftime("%B %Y")
        date_para = date_frame.paragraphs[0]
        date_para.font.size = Pt(18)
        date_para.font.color.rgb = self.COLORS['white']
    
    def add_section_divider(self, title: str, subtitle: str, icon: str = ""):
        """Add beautiful section divider"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide, 'section')
        
        # Large icon
        if icon:
            icon_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(2.5), Inches(1.5))
            icon_frame = icon_box.text_frame
            icon_frame.text = icon
            icon_para = icon_frame.paragraphs[0]
            icon_para.font.size = Pt(100)
            icon_para.alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(2), Inches(3.8), Inches(9.333), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(3), Inches(5), Inches(7.333), Inches(0.6))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(20)
            subtitle_para.font.color.rgb = self.COLORS['white']
            subtitle_para.alignment = PP_ALIGN.CENTER
    
    def add_executive_summary(self, summary: str, highlights: List[str]):
        """Add executive summary slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide)
        
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(0.5), Inches(12), Inches(0.7)
        )
        fill = header.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['primary']
        header.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.6), Inches(11), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "📋 Executive Summary"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        
        # Summary text
        summary_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(1.8))
        summary_frame = summary_box.text_frame
        summary_frame.word_wrap = True
        summary_frame.text = summary
        summary_para = summary_frame.paragraphs[0]
        summary_para.font.size = Pt(18)
        summary_para.font.color.rgb = self.COLORS['text']
        summary_para.line_spacing = 1.4
        
        # Key highlights in cards
        highlight_label = slide.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(12), Inches(0.4))
        hl_frame = highlight_label.text_frame
        hl_frame.text = "Key Highlights"
        hl_para = hl_frame.paragraphs[0]
        hl_para.font.size = Pt(22)
        hl_para.font.bold = True
        hl_para.font.color.rgb = self.COLORS['primary']
        
        colors = [self.COLORS['secondary'], self.COLORS['success'], self.COLORS['accent'], self.COLORS['purple']]
        icons = ['🎯', '💡', '🚀', '⭐']
        
        for i, highlight in enumerate(highlights[:4]):
            col = i % 2
            row = i // 2
            left = 0.7 + (col * 6.2)
            top = 4.2 + (row * 1.4)
            
            # Card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top), Inches(5.8), Inches(1.2)
            )
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['white']
            card.line.color.rgb = colors[i]
            card.line.width = Pt(2)
            card.shadow.inherit = False
            
            # Icon circle
            icon_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left + 0.15), Inches(top + 0.25), Inches(0.7), Inches(0.7)
            )
            fill = icon_circle.fill
            fill.solid()
            fill.fore_color.rgb = colors[i]
            icon_circle.line.fill.background()
            
            icon_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.25), Inches(0.7), Inches(0.7))
            icon_frame = icon_box.text_frame
            icon_frame.text = icons[i]
            icon_para = icon_frame.paragraphs[0]
            icon_para.font.size = Pt(28)
            icon_para.alignment = PP_ALIGN.CENTER
            icon_para.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Text
            text_box = slide.shapes.add_textbox(Inches(left + 1), Inches(top + 0.15), Inches(4.6), Inches(0.9))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame.text = highlight
            text_para = text_frame.paragraphs[0]
            text_para.font.size = Pt(14)
            text_para.font.color.rgb = self.COLORS['text']
    
    def add_value_proposition_slide(self, propositions: List[str]):
        """Add value proposition slide with impressive layout"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide)
        
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(0.5), Inches(12), Inches(0.7)
        )
        fill = header.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['secondary']
        header.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.6), Inches(11), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "💎 Value Propositions"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        
        # Value cards in staggered layout
        card_colors = [self.COLORS['secondary'], self.COLORS['accent'], self.COLORS['success']]
        
        for i, prop in enumerate(propositions[:3]):
            left_offset = 0.5 + (i * 0.3)
            top = 2 + (i * 1.6)
            
            # Main card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left_offset), Inches(top), Inches(11.5), Inches(1.4)
            )
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = card_colors[i]
            card.line.fill.background()
            card.shadow.inherit = False
            
            # Number badge
            badge = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left_offset + 0.3), Inches(top + 0.35), Inches(0.7), Inches(0.7)
            )
            fill = badge.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['white']
            badge.line.fill.background()
            
            num_box = slide.shapes.add_textbox(Inches(left_offset + 0.3), Inches(top + 0.35), Inches(0.7), Inches(0.7))
            num_frame = num_box.text_frame
            num_frame.text = str(i + 1)
            num_para = num_frame.paragraphs[0]
            num_para.font.size = Pt(24)
            num_para.font.bold = True
            num_para.font.color.rgb = card_colors[i]
            num_para.alignment = PP_ALIGN.CENTER
            num_para.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Value text
            text_box = slide.shapes.add_textbox(Inches(left_offset + 1.2), Inches(top + 0.25), Inches(10), Inches(0.9))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame.text = prop
            text_para = text_frame.paragraphs[0]
            text_para.font.size = Pt(18)
            text_para.font.color.rgb = self.COLORS['white']
            text_para.font.bold = True
    
    def add_architecture_overview(self, pattern: str, components: List[Dict], tech_stack: List[str]):
        """Add technical architecture overview"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide)
        
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(0.5), Inches(12), Inches(0.7)
        )
        fill = header.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['purple']
        header.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.6), Inches(11), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "🏗️ Technical Architecture"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        
        # Architecture pattern badge
        pattern_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(1.5), Inches(11.333), Inches(0.6)
        )
        fill = pattern_box.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['light']
        pattern_box.line.color.rgb = self.COLORS['purple']
        pattern_box.line.width = Pt(2)
        
        pattern_text = slide.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(11), Inches(0.4))
        pattern_frame = pattern_text.text_frame
        pattern_frame.text = f"Architecture Pattern: {pattern}"
        pattern_para = pattern_frame.paragraphs[0]
        pattern_para.font.size = Pt(20)
        pattern_para.font.bold = True
        pattern_para.font.color.rgb = self.COLORS['purple']
        
        # Component boxes in grid
        type_colors = {
            'frontend': self.COLORS['secondary'],
            'backend': self.COLORS['success'],
            'database': self.COLORS['accent'],
            'service': self.COLORS['purple'],
            'default': self.COLORS['primary']
        }
        
        cols = 3
        box_width = 3.6
        box_height = 1.6
        spacing = 0.4
        
        for i, comp in enumerate(components[:6]):
            row = i // cols
            col = i % cols
            left = 1 + (col * (box_width + spacing))
            top = 2.5 + (row * (box_height + spacing))
            
            comp_type = comp.get('type', 'default').lower()
            comp_color = type_colors.get(comp_type, type_colors['default'])
            
            # Component card with gradient effect
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top), Inches(box_width), Inches(box_height)
            )
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = comp_color
            card.line.fill.background()
            card.shadow.inherit = False
            
            # Type badge
            type_badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left + 0.15), Inches(top + 0.15), Inches(1.5), Inches(0.35)
            )
            fill = type_badge.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['white']
            type_badge.line.fill.background()
            
            type_text = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.18), Inches(1.4), Inches(0.29))
            type_frame = type_text.text_frame
            type_frame.text = comp_type.upper()
            type_para = type_frame.paragraphs[0]
            type_para.font.size = Pt(9)
            type_para.font.bold = True
            type_para.font.color.rgb = comp_color
            type_para.alignment = PP_ALIGN.CENTER
            
            # Component name
            name_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.6), Inches(box_width - 0.4), Inches(0.4))
            name_frame = name_box.text_frame
            name_frame.text = comp.get('name', 'Component')
            name_para = name_frame.paragraphs[0]
            name_para.font.size = Pt(18)
            name_para.font.bold = True
            name_para.font.color.rgb = self.COLORS['white']
            name_para.alignment = PP_ALIGN.CENTER
            
            # Responsibility
            resp_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 1.05), Inches(box_width - 0.3), Inches(0.45))
            resp_frame = resp_box.text_frame
            resp_frame.word_wrap = True
            resp_text = comp.get('responsibility', 'N/A')
            if len(resp_text) > 60:
                resp_text = resp_text[:57] + "..."
            resp_frame.text = resp_text
            resp_para = resp_frame.paragraphs[0]
            resp_para.font.size = Pt(11)
            resp_para.font.color.rgb = self.COLORS['white']
            resp_para.alignment = PP_ALIGN.CENTER
        
        # Tech stack at bottom
        if tech_stack:
            tech_label = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(11.333), Inches(0.3))
            tech_frame = tech_label.text_frame
            tech_frame.text = "Technology Stack: " + " • ".join(tech_stack[:8])
            tech_para = tech_frame.paragraphs[0]
            tech_para.font.size = Pt(14)
            tech_para.font.color.rgb = self.COLORS['text_light']
            tech_para.alignment = PP_ALIGN.CENTER
    
    def add_roadmap_slide(self, phases: List[Dict]):
        """Add implementation roadmap with timeline"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide)
        
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(0.5), Inches(12), Inches(0.7)
        )
        fill = header.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['accent']
        header.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.6), Inches(11), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "🗓️ Implementation Roadmap"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        
        # Timeline
        timeline_colors = [self.COLORS['secondary'], self.COLORS['success'], self.COLORS['accent'], self.COLORS['purple']]
        
        for i, phase in enumerate(phases[:4]):
            top = 2 + (i * 1.3)
            
            # Timeline dot and line
            if i > 0:
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(1.6), Inches(top - 0.65), Inches(0.08), Inches(0.65)
                )
                fill = line.fill
                fill.solid()
                fill.fore_color.rgb = timeline_colors[i-1]
                line.line.fill.background()
            
            # Phase dot
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(1.4), Inches(top), Inches(0.5), Inches(0.5)
            )
            fill = dot.fill
            fill.solid()
            fill.fore_color.rgb = timeline_colors[i]
            dot.line.color.rgb = self.COLORS['white']
            dot.line.width = Pt(3)
            
            # Phase card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(2.2), Inches(top), Inches(10.2), Inches(1.1)
            )
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['white']
            card.line.color.rgb = timeline_colors[i]
            card.line.width = Pt(2)
            card.shadow.inherit = False
            
            # Phase header
            phase_header = slide.shapes.add_textbox(Inches(2.5), Inches(top + 0.15), Inches(5), Inches(0.35))
            phase_frame = phase_header.text_frame
            phase_frame.text = f"{phase.get('phase', f'Phase {i+1}')}: {phase.get('title', 'Implementation')}"
            phase_para = phase_frame.paragraphs[0]
            phase_para.font.size = Pt(16)
            phase_para.font.bold = True
            phase_para.font.color.rgb = timeline_colors[i]
            
            # Duration badge
            duration_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(8), Inches(top + 0.12), Inches(1.8), Inches(0.38)
            )
            fill = duration_box.fill
            fill.solid()
            fill.fore_color.rgb = timeline_colors[i]
            duration_box.line.fill.background()
            
            dur_text = slide.shapes.add_textbox(Inches(8.1), Inches(top + 0.17), Inches(1.6), Inches(0.28))
            dur_frame = dur_text.text_frame
            dur_frame.text = f"⏱️ {phase.get('duration', 'TBD')}"
            dur_para = dur_frame.paragraphs[0]
            dur_para.font.size = Pt(12)
            dur_para.font.bold = True
            dur_para.font.color.rgb = self.COLORS['white']
            dur_para.alignment = PP_ALIGN.CENTER
            
            # Deliverables
            deliv_box = slide.shapes.add_textbox(Inches(2.5), Inches(top + 0.55), Inches(9.5), Inches(0.45))
            deliv_frame = deliv_box.text_frame
            deliv_frame.word_wrap = True
            deliverables = phase.get('deliverables', [])
            deliv_text = " • ".join(deliverables[:3]) if deliverables else "Key deliverables"
            if len(deliv_text) > 100:
                deliv_text = deliv_text[:97] + "..."
            deliv_frame.text = deliv_text
            deliv_para = deliv_frame.paragraphs[0]
            deliv_para.font.size = Pt(12)
            deliv_para.font.color.rgb = self.COLORS['text']
    
    def add_feature_matrix(self, features: List[str], capabilities: List[str]):
        """Add features and capabilities matrix"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide)
        
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(0.5), Inches(12), Inches(0.7)
        )
        fill = header.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['success']
        header.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.6), Inches(11), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "⚡ Features & Capabilities"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        
        # Two column layout
        # Left: Features
        feature_label = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.5)
        )
        fill = feature_label.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['secondary']
        feature_label.line.fill.background()
        
        fl_text = slide.shapes.add_textbox(Inches(0.9), Inches(1.58), Inches(5.4), Inches(0.34))
        fl_frame = fl_text.text_frame
        fl_frame.text = "🎯 Core Features"
        fl_para = fl_frame.paragraphs[0]
        fl_para.font.size = Pt(20)
        fl_para.font.bold = True
        fl_para.font.color.rgb = self.COLORS['white']
        
        for i, feature in enumerate(features[:5]):
            top = 2.2 + (i * 0.85)
            
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.7), Inches(top), Inches(5.8), Inches(0.7)
            )
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['white']
            card.line.color.rgb = self.COLORS['border']
            card.line.width = Pt(1)
            
            # Checkmark
            check_box = slide.shapes.add_textbox(Inches(0.95), Inches(top + 0.15), Inches(0.4), Inches(0.4))
            check_frame = check_box.text_frame
            check_frame.text = "✓"
            check_para = check_frame.paragraphs[0]
            check_para.font.size = Pt(24)
            check_para.font.bold = True
            check_para.font.color.rgb = self.COLORS['success']
            check_para.alignment = PP_ALIGN.CENTER
            
            text_box = slide.shapes.add_textbox(Inches(1.5), Inches(top + 0.15), Inches(4.8), Inches(0.4))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame.text = feature
            text_para = text_frame.paragraphs[0]
            text_para.font.size = Pt(14)
            text_para.font.color.rgb = self.COLORS['text']
        
        # Right: Capabilities
        cap_label = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5)
        )
        fill = cap_label.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['accent']
        cap_label.line.fill.background()
        
        cl_text = slide.shapes.add_textbox(Inches(7), Inches(1.58), Inches(5.4), Inches(0.34))
        cl_frame = cl_text.text_frame
        cl_frame.text = "💪 Technical Capabilities"
        cl_para = cl_frame.paragraphs[0]
        cl_para.font.size = Pt(20)
        cl_para.font.bold = True
        cl_para.font.color.rgb = self.COLORS['white']
        
        for i, capability in enumerate(capabilities[:5]):
            top = 2.2 + (i * 0.85)
            
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.8), Inches(top), Inches(5.8), Inches(0.7)
            )
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['white']
            card.line.color.rgb = self.COLORS['border']
            card.line.width = Pt(1)
            
            # Star icon
            star_box = slide.shapes.add_textbox(Inches(7.05), Inches(top + 0.15), Inches(0.4), Inches(0.4))
            star_frame = star_box.text_frame
            star_frame.text = "⭐"
            star_para = star_frame.paragraphs[0]
            star_para.font.size = Pt(20)
            star_para.alignment = PP_ALIGN.CENTER
            
            text_box = slide.shapes.add_textbox(Inches(7.6), Inches(top + 0.15), Inches(4.8), Inches(0.4))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame.text = capability
            text_para = text_frame.paragraphs[0]
            text_para.font.size = Pt(14)
            text_para.font.color.rgb = self.COLORS['text']
    
    def add_security_compliance(self, security: List[str], compliance: List[str]):
        """Add security and compliance slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide)
        
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(0.5), Inches(12), Inches(0.7)
        )
        fill = header.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['primary']
        header.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.6), Inches(11), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "🔒 Security & Compliance"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        
        # Security shield icon
        shield = slide.shapes.add_textbox(Inches(5.8), Inches(1.5), Inches(1.7), Inches(1))
        shield_frame = shield.text_frame
        shield_frame.text = "🛡️"
        shield_para = shield_frame.paragraphs[0]
        shield_para.font.size = Pt(80)
        shield_para.alignment = PP_ALIGN.CENTER
        
        # Security measures
        sec_top = 2.8
        for i, measure in enumerate(security[:4]):
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1), Inches(sec_top + (i * 0.75)), Inches(5.5), Inches(0.6)
            )
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['white']
            card.line.color.rgb = self.COLORS['primary']
            card.line.width = Pt(2)
            
            text_box = slide.shapes.add_textbox(Inches(1.3), Inches(sec_top + (i * 0.75) + 0.1), Inches(5), Inches(0.4))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame.text = f"🔐 {measure}"
            text_para = text_frame.paragraphs[0]
            text_para.font.size = Pt(13)
            text_para.font.color.rgb = self.COLORS['text']
        
        # Compliance badges
        comp_top = 2.8
        for i, comp in enumerate(compliance[:4]):
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.9), Inches(comp_top + (i * 0.75)), Inches(5.5), Inches(0.6)
            )
            fill = badge.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['success']
            badge.line.fill.background()
            
            text_box = slide.shapes.add_textbox(Inches(7.2), Inches(comp_top + (i * 0.75) + 0.1), Inches(5), Inches(0.4))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame.text = f"✓ {comp}"
            text_para = text_frame.paragraphs[0]
            text_para.font.size = Pt(13)
            text_para.font.bold = True
            text_para.font.color.rgb = self.COLORS['white']
    
    def add_investment_summary(self, resources: Dict):
        """Add investment and resource summary"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide)
        
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(0.5), Inches(12), Inches(0.7)
        )
        fill = header.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['gold']
        header.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.6), Inches(11), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "💼 Investment Overview"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        
        # Investment cards
        cards_data = [
            ("Infrastructure", resources.get('infrastructure_needs', []), self.COLORS['secondary'], "🖥️"),
            ("Team Requirements", resources.get('team_requirements', []), self.COLORS['success'], "👥"),
            ("Timeline", [resources.get('estimated_timeline', 'TBD')], self.COLORS['accent'], "⏱️"),
        ]
        
        for i, (title, items, color, icon) in enumerate(cards_data):
            left = 1 + (i * 3.9)
            
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(1.8), Inches(3.6), Inches(4.2)
            )
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['white']
            card.line.color.rgb = color
            card.line.width = Pt(3)
            card.shadow.inherit = False
            
            # Icon header
            icon_header = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(1.8), Inches(3.6), Inches(0.8)
            )
            fill = icon_header.fill
            fill.solid()
            fill.fore_color.rgb = color
            icon_header.line.fill.background()
            
            icon_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(1.9), Inches(0.6), Inches(0.6))
            icon_frame = icon_box.text_frame
            icon_frame.text = icon
            icon_para = icon_frame.paragraphs[0]
            icon_para.font.size = Pt(32)
            
            title_box = slide.shapes.add_textbox(Inches(left + 0.9), Inches(2.05), Inches(2.5), Inches(0.5))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(18)
            title_para.font.bold = True
            title_para.font.color.rgb = self.COLORS['white']
            
            # Items
            item_top = 2.8
            for j, item in enumerate(items[:4]):
                item_text = slide.shapes.add_textbox(Inches(left + 0.2), Inches(item_top + (j * 0.7)), Inches(3.2), Inches(0.6))
                item_frame = item_text.text_frame
                item_frame.word_wrap = True
                item_text_str = str(item)
                if len(item_text_str) > 50:
                    item_text_str = item_text_str[:47] + "..."
                item_frame.text = f"• {item_text_str}"
                item_para = item_frame.paragraphs[0]
                item_para.font.size = Pt(12)
                item_para.font.color.rgb = self.COLORS['text']
                item_para.line_spacing = 1.2
    
    def add_closing_slide(self, title: str = "Thank You"):
        """Add professional closing slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide, 'section')
        
        # Thank you text
        thank_box = slide.shapes.add_textbox(Inches(3), Inches(2.5), Inches(7.333), Inches(1))
        thank_frame = thank_box.text_frame
        thank_frame.text = title
        thank_para = thank_frame.paragraphs[0]
        thank_para.font.size = Pt(56)
        thank_para.font.bold = True
        thank_para.font.color.rgb = self.COLORS['white']
        thank_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(3), Inches(3.8), Inches(7.333), Inches(0.6))
        sub_frame = sub_box.text_frame
        sub_frame.text = "Ready to Transform Your Business"
        sub_para = sub_frame.paragraphs[0]
        sub_para.font.size = Pt(24)
        sub_para.font.color.rgb = self.COLORS['white']
        sub_para.alignment = PP_ALIGN.CENTER
        
        # Contact info
        contact_box = slide.shapes.add_textbox(Inches(3), Inches(5), Inches(7.333), Inches(0.8))
        contact_frame = contact_box.text_frame
        contact_frame.text = f"{self.company_name}\nLet's discuss how we can help you succeed"
        contact_para = contact_frame.paragraphs[0]
        contact_para.font.size = Pt(16)
        contact_para.font.color.rgb = self.COLORS['white']
        contact_para.alignment = PP_ALIGN.CENTER
        contact_para.line_spacing = 1.5
    
    def save_to_bytes(self) -> bytes:
        """Save presentation to bytes"""
        buffer = io.BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()


def generate_business_presentation(analysis_data: Dict, company_name: str) -> bytes:
    """Generate stunning business proposal presentation"""
    ppt = BusinessPPTGenerator(company_name)
    
    # 1. Cover Slide
    ppt.add_cover_slide(
        "Architecture & Technical Proposal",
        "Comprehensive Solution Analysis"
    )
    
    # 2. Executive Summary
    business = analysis_data.get('business', {})
    if business:
        ppt.add_section_divider("Executive Overview", "Business Value & Strategic Impact", "📊")
        
        exec_summary = business.get('executive_summary', 'Comprehensive technical solution designed for optimal performance and scalability.')
        highlights = [
            business.get('business_problem', 'Solves critical business challenges'),
            f"Target Users: {', '.join(business.get('target_users', ['Enterprise clients'])[:2])}",
            f"Key Advantage: {business.get('competitive_advantages', ['Advanced technology'])[0] if business.get('competitive_advantages') else 'Cutting-edge solution'}",
            f"Expected ROI: {business.get('business_benefits', ['Significant value'])[0] if business.get('business_benefits') else 'High return on investment'}"
        ]
        ppt.add_executive_summary(exec_summary, highlights)
    
    # 3. Value Propositions
    if business and business.get('value_propositions'):
        ppt.add_value_proposition_slide(business['value_propositions'])
    
    # 4. Technical Architecture
    arch = analysis_data.get('architecture', {})
    if arch:
        ppt.add_section_divider("Technical Solution", "Architecture & Implementation Details", "🏗️")
        
        ppt.add_architecture_overview(
            arch.get('architecture_pattern', 'Modern Architecture'),
            arch.get('main_components', []),
            arch.get('technology_stack', [])
        )
    
    # 5. Features & Capabilities
    capabilities = analysis_data.get('capabilities', {})
    if capabilities:
        ppt.add_feature_matrix(
            capabilities.get('core_features', []),
            capabilities.get('technical_capabilities', [])
        )
    
    # 6. Implementation Roadmap
    roadmap = analysis_data.get('roadmap', {})
    if roadmap and roadmap.get('phases'):
        ppt.add_section_divider("Implementation Plan", "Phased Approach to Success", "🗓️")
        ppt.add_roadmap_slide(roadmap['phases'])
    
    # 7. Security & Compliance
    security = analysis_data.get('security', {})
    if security:
        ppt.add_security_compliance(
            security.get('security_measures', []),
            security.get('compliance_standards', ['Industry Best Practices', 'Data Protection', 'Secure Development'])
        )
    
    # 8. Investment Overview
    resources = analysis_data.get('resources', {})
    if resources:
        ppt.add_investment_summary(resources)
    
    # 9. Closing Slide
    ppt.add_closing_slide()
    
    return ppt.save_to_bytes()


# Streamlit UI
st.title("📊 Business Proposal & Architecture Generator")
st.markdown("Transform your codebase into a **stunning business proposal** with AI-powered analysis")

# Sidebar
with st.sidebar:
    st.header("⚙️ Configuration")
    
    company_name = st.text_input(
        "Company/Project Name",
        value="Acme Corporation",
        help="Name to display on the presentation"
    )
    
    st.markdown("---")
    st.subheader("Azure OpenAI Settings")
    
    azure_endpoint = st.text_input(
        "Azure Endpoint",
        placeholder="https://your-resource.openai.azure.com/",
        type="password"
    )
    
    azure_key = st.text_input(
        "Azure API Key",
        type="password"
    )
    
    azure_deployment = st.text_input(
        "Deployment Name",
        placeholder="gpt-4"
    )
    
    st.markdown("---")
    st.markdown("### 🎯 What You'll Get")
    st.markdown("""
    - **Executive Summary** with business value
    - **Value Propositions** for stakeholders
    - **Technical Architecture** diagrams
    - **Implementation Roadmap** with phases
    - **Security & Compliance** overview
    - **Investment Analysis** and resources
    """)

# Main content
st.markdown("---")

col1, col2 = st.columns([2, 1])

with col1:
    uploaded_file = st.file_uploader(
        "📦 Upload Codebase (ZIP)", 
        type=['zip'],
        help="Upload a ZIP file containing your source code"
    )

with col2:
    st.info("✨ Generates professional business proposals from code")

if uploaded_file and azure_endpoint and azure_key and azure_deployment:
    
    if st.button("🚀 Generate Business Proposal", type="primary", use_container_width=True):
        
        with st.spinner("Analyzing and generating presentation..."):
            try:
                with tempfile.TemporaryDirectory() as temp_dir:
                    # Save uploaded file
                    zip_path = os.path.join(temp_dir, "codebase.zip")
                    with open(zip_path, 'wb') as f:
                        f.write(uploaded_file.getvalue())
                    
                    extract_dir = os.path.join(temp_dir, "extracted")
                    os.makedirs(extract_dir, exist_ok=True)
                    
                    # Initialize analyzer
                    analyzer = CodebaseAnalyzer(azure_endpoint, azure_key, azure_deployment)
                    
                    # Progress tracking
                    progress_bar = st.progress(0)
                    status_text = st.empty()
                    
                    status_text.text("📂 Extracting codebase...")
                    progress_bar.progress(10)
                    
                    file_contents = analyzer.extract_files(zip_path, extract_dir)
                    
                    st.success(f"✅ Extracted {len(file_contents)} files")
                    progress_bar.progress(20)
                    
                    # Analyze codebase
                    def update_progress(message):
                        status_text.text(f"🔍 {message}")
                    
                    status_text.text("🔍 Performing AI analysis...")
                    analysis_data = analyzer.analyze_codebase(file_contents, update_progress)
                    
                    progress_bar.progress(80)
                    status_text.text("✨ Generating stunning presentation...")
                    
                    # Generate presentation
                    ppt_bytes = generate_business_presentation(analysis_data, company_name)
                    
                    progress_bar.progress(100)
                    status_text.text("✅ Presentation ready!")
                    
                    # Store in session state
                    st.session_state.analysis_complete = True
                    st.session_state.analysis_data = analysis_data
                    st.session_state.ppt_bytes = ppt_bytes
                    
                    st.balloons()
                    
            except Exception as e:
                st.error(f"❌ Error: {str(e)}")
                st.exception(e)

# Display results
if st.session_state.analysis_complete:
    st.success("🎉 Business Proposal Generated Successfully!")
    
    # Download button - prominent
    st.markdown("---")
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        st.download_button(
            label="📥 Download Business Proposal (PowerPoint)",
            data=st.session_state.ppt_bytes,
            file_name=f"business_proposal_{datetime.now().strftime('%Y%m%d')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
            type="primary"
        )
    
    # Analysis preview
    st.markdown("---")
    st.subheader("📊 Analysis Preview")
    
    analysis = st.session_state.analysis_data
    
    # Tabs for different sections
    tabs = st.tabs([
        "💼 Business Value", 
        "🏗️ Architecture", 
        "⚡ Capabilities", 
        "🗓️ Roadmap",
        "🔒 Security",
        "💰 Investment"
    ])
    
    # Business Value Tab
    with tabs[0]:
        business = analysis.get('business', {})
        if business:
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("#### 🎯 Business Problem")
                st.info(business.get('business_problem', 'N/A'))
                
                st.markdown("#### 👥 Target Users")
                for user in business.get('target_users', []):
                    st.markdown(f"- {user}")
            
            with col2:
                st.markdown("#### 💎 Value Propositions")
                for i, prop in enumerate(business.get('value_propositions', []), 1):
                    st.markdown(f"**{i}.** {prop}")
                
                st.markdown("#### 🏆 Competitive Advantages")
                for adv in business.get('competitive_advantages', []):
                    st.markdown(f"- {adv}")
            
            st.markdown("#### 📈 Business Benefits")
            benefits_cols = st.columns(3)
            for i, benefit in enumerate(business.get('business_benefits', [])):
                with benefits_cols[i % 3]:
                    st.success(benefit)
    
    # Architecture Tab
    with tabs[1]:
        arch = analysis.get('architecture', {})
        if arch:
            col1, col2 = st.columns([1, 1])
            
            with col1:
                st.markdown("#### 🏗️ Architecture Pattern")
                st.code(arch.get('architecture_pattern', 'N/A'), language=None)
                
                st.markdown("#### 📦 Technology Stack")
                tech_stack = arch.get('technology_stack', [])
                for tech in tech_stack:
                    st.markdown(f"- `{tech}`")
            
            with col2:
                st.markdown("#### 🔧 Main Components")
                components = arch.get('main_components', [])
                for comp in components:
                    with st.expander(f"**{comp.get('name', 'Component')}** ({comp.get('type', 'N/A')})"):
                        st.write(comp.get('responsibility', 'N/A'))
            
            if arch.get('scalability'):
                st.markdown("#### 📊 Scalability")
                st.info(arch.get('scalability'))
    
    # Capabilities Tab
    with tabs[2]:
        capabilities = analysis.get('capabilities', {})
        if capabilities:
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("#### 🎯 Core Features")
                for feature in capabilities.get('core_features', []):
                    st.markdown(f"✓ {feature}")
            
            with col2:
                st.markdown("#### 💪 Technical Capabilities")
                for cap in capabilities.get('technical_capabilities', []):
                    st.markdown(f"⭐ {cap}")
            
            st.markdown("#### ⚡ Performance Characteristics")
            perf_cols = st.columns(3)
            for i, perf in enumerate(capabilities.get('performance_characteristics', [])):
                with perf_cols[i % 3]:
                    st.info(perf)
            
            if capabilities.get('integration_points'):
                st.markdown("#### 🔌 Integration Points")
                for integration in capabilities.get('integration_points', []):
                    st.markdown(f"- {integration}")
    
    # Roadmap Tab
    with tabs[3]:
        roadmap = analysis.get('roadmap', {})
        if roadmap:
            phases = roadmap.get('phases', [])
            for i, phase in enumerate(phases, 1):
                with st.expander(f"**{phase.get('phase', f'Phase {i}')}**: {phase.get('title', 'Implementation')}", expanded=(i==1)):
                    col1, col2 = st.columns([2, 1])
                    
                    with col1:
                        st.markdown("**Deliverables:**")
                        for deliv in phase.get('deliverables', []):
                            st.markdown(f"- {deliv}")
                    
                    with col2:
                        st.metric("Duration", phase.get('duration', 'TBD'))
            
            if roadmap.get('milestones'):
                st.markdown("#### 🎯 Key Milestones")
                milestone_cols = st.columns(2)
                for i, milestone in enumerate(roadmap.get('milestones', [])):
                    with milestone_cols[i % 2]:
                        st.success(f"✓ {milestone}")
            
            if roadmap.get('success_criteria'):
                st.markdown("#### 📊 Success Criteria")
                for criteria in roadmap.get('success_criteria', []):
                    st.markdown(f"- {criteria}")
    
    # Security Tab
    with tabs[4]:
        security = analysis.get('security', {})
        if security:
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("#### 🛡️ Security Measures")
                for measure in security.get('security_measures', []):
                    st.markdown(f"🔐 {measure}")
                
                st.markdown("#### 🔒 Authentication Methods")
                for auth in security.get('authentication_methods', []):
                    st.markdown(f"- {auth}")
            
            with col2:
                st.markdown("#### ✓ Compliance Standards")
                for comp in security.get('compliance_standards', []):
                    st.success(comp)
                
                st.markdown("#### 🗄️ Data Protection")
                for data in security.get('data_protection', []):
                    st.markdown(f"- {data}")
    
    # Investment Tab
    with tabs[5]:
        resources = analysis.get('resources', {})
        if resources:
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.markdown("#### 🖥️ Infrastructure Needs")
                for infra in resources.get('infrastructure_needs', []):
                    st.markdown(f"- {infra}")
            
            with col2:
                st.markdown("#### 👥 Team Requirements")
                for team in resources.get('team_requirements', []):
                    st.markdown(f"- {team}")
            
            with col3:
                st.markdown("#### ⏱️ Timeline")
                st.metric("Estimated", resources.get('estimated_timeline', 'TBD'))
            
            if resources.get('maintenance_considerations'):
                st.markdown("#### 🔧 Maintenance Considerations")
                maint_cols = st.columns(2)
                for i, maint in enumerate(resources.get('maintenance_considerations', [])):
                    with maint_cols[i % 2]:
                        st.info(maint)
    
    # Reset button
    st.markdown("---")
    if st.button("🔄 Analyze Another Codebase", use_container_width=True):
        st.session_state.analysis_complete = False
        st.session_state.analysis_data = None
        st.session_state.ppt_bytes = None
        st.rerun()

elif not (azure_endpoint and azure_key and azure_deployment):
    st.warning("⚠️ Please configure Azure OpenAI settings in the sidebar to get started.")
    
    # Show features
    st.markdown("---")
    st.markdown("### 🌟 What Makes This Special?")
    
    feat_cols = st.columns(3)
    
    with feat_cols[0]:
        st.markdown("""
        #### 💼 Business-Focused
        - Executive summaries
        - Value propositions
        - ROI analysis
        - Stakeholder insights
        """)
    
    with feat_cols[1]:
        st.markdown("""
        #### 🎨 Stunning Design
        - Modern 16:9 layouts
        - Professional color schemes
        - Visual diagrams
        - Premium aesthetics
        """)
    
    with feat_cols[2]:
        st.markdown("""
        #### 🤖 AI-Powered
        - Intelligent analysis
        - Architecture detection
        - Security assessment
        - Roadmap generation
        """)
    
    st.markdown("---")
    st.info("💡 **Tip:** This tool is perfect for creating proposals for clients, investors, or internal stakeholders!")

else:
    st.info("👆 Upload a ZIP file containing your codebase to begin generating your business proposal.")
    
    # Example use cases
    with st.expander("📚 Example Use Cases"):
        st.markdown("""
        **Perfect for:**
        - 🤝 Client presentations and proposals
        - 💰 Investor pitch decks (technical section)
        - 📊 Internal architecture reviews
        - 🎯 Stakeholder communications
        - 📖 Technical documentation
        - 🔄 System modernization proposals
        - 🌐 RFP responses
        - 🚀 Product launch materials
        """)
    
    with st.expander("🎯 Best Practices"):
        st.markdown("""
        **For Best Results:**
        - Include README and documentation files
        - Ensure main code files are present
        - Remove unnecessary build artifacts
        - Include configuration files (package.json, requirements.txt, etc.)
        - Typical ZIP size: 1-50MB
        - Supported languages: Python, JavaScript, Java, C++, and more
        """)

