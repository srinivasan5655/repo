import streamlit as st
import zipfile
import os
import tempfile
from pathlib import Path
from openai import AzureOpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Inches as DocInches, Pt as DocPt, RGBColor as DocRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
import io
import json
from typing import Dict, List, Tuple
import re
from datetime import datetime

# Page config
st.set_page_config(
    page_title="Business Proposal & Architecture Generator",
    page_icon="📊",
    layout="wide"
)

# Initialize session state
if 'analysis_complete' not in st.session_state:
    st.session_state.analysis_complete = False
if 'analysis_data' not in st.session_state:
    st.session_state.analysis_data = None
if 'input_mode' not in st.session_state:
    st.session_state.input_mode = 'upload'

class CodebaseAnalyzer:
    """Analyzes codebase and generates architecture insights"""
    
    def __init__(self, azure_endpoint: str, azure_key: str, azure_deployment: str):
        self.client = AzureOpenAI(
            azure_endpoint=azure_endpoint,
            api_key=azure_key,
            api_version="2024-05-01-preview"
        )
        self.deployment = azure_deployment
        
    def extract_files(self, zip_path: str, extract_to: str) -> Dict[str, str]:
        """Extract and read files from zip"""
        file_contents = {}
        
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(extract_to)
            
        # Read all code files
        for root, dirs, files in os.walk(extract_to):
            dirs[:] = [d for d in dirs if d not in ['.git', 'node_modules', '__pycache__', 'venv', '.venv', 'dist', 'build']]
            
            for file in files:
                if self._is_code_file(file):
                    file_path = os.path.join(root, file)
                    try:
                        with open(file_path, 'r', encoding='utf-8') as f:
                            relative_path = os.path.relpath(file_path, extract_to)
                            file_contents[relative_path] = f.read()
                    except Exception as e:
                        st.warning(f"Could not read {file}: {e}")
                        
        return file_contents
    
    def _is_code_file(self, filename: str) -> bool:
        """Check if file is a code file"""
        code_extensions = [
            '.py', '.js', '.jsx', '.ts', '.tsx', '.java', '.cpp', '.c', '.h',
            '.cs', '.go', '.rs', '.rb', '.php', '.swift', '.kt', '.scala',
            '.json', '.yaml', '.yml', '.xml', '.sql', '.sh', '.md', '.txt'
        ]
        return any(filename.endswith(ext) for ext in code_extensions)
    
    def analyze_from_text(self, user_input: str, progress_callback=None) -> Dict:
        """Analyze requirements from user text input"""
        analyses = {}
        
        # 1. Business Value & Executive Summary
        if progress_callback:
            progress_callback("Analyzing business requirements...")
        
        business_prompt = f"""Analyze this project description from a BUSINESS PERSPECTIVE and provide:
1. Main business problem it solves
2. Key value propositions
3. Target users/stakeholders
4. Business benefits and ROI potential
5. Competitive advantages

Project Description:
{user_input}

strictly Respond in JSON format without any marking or explanations:
{{
    "business_problem": "...",
    "value_propositions": ["..."],
    "target_users": ["..."],
    "business_benefits": ["..."],
    "competitive_advantages": ["..."],
    "executive_summary": "..."
}}"""
        
        analyses['business'] = self._call_llm(business_prompt)
        
        # 2. Application Architecture
        if progress_callback:
            progress_callback("Designing application architecture...")
        
        app_arch_prompt = f"""Design detailed application architecture:
1. Architecture layers and tiers
2. Application components and services
3. Communication patterns
4. Deployment architecture
5. Scalability and performance design

Project Description:
{user_input}

strictly Respond in JSON format without any marking or explanations:
{{
    "architecture_layers": [{{"name": "...", "description": "...", "components": ["..."]}}],
    "services": [{{"name": "...", "type": "...", "responsibility": "...", "technologies": ["..."]}}],
    "communication_patterns": ["..."],
    "deployment_model": "...",
    "scalability_approach": "...",
    "performance_optimization": ["..."]
}}"""
        
        analyses['app_architecture'] = self._call_llm(app_arch_prompt)
        
        # 3. Technical Architecture
        if progress_callback:
            progress_callback("Designing technical architecture...")
        
        tech_arch_prompt = f"""Design comprehensive technical architecture:
1. Technology stack (frontend, backend, database, infrastructure)
2. Integration architecture
3. Security architecture
4. DevOps and CI/CD pipeline
5. Monitoring and logging

Project Description:
{user_input}

strictly Respond in JSON format without any marking or explanations:
{{
    "frontend_stack": [{{"technology": "...", "purpose": "..."}}],
    "backend_stack": [{{"technology": "...", "purpose": "..."}}],
    "database_stack": [{{"technology": "...", "purpose": "..."}}],
    "infrastructure": [{{"component": "...", "technology": "...", "purpose": "..."}}],
    "integration_points": [{{"system": "...", "method": "...", "protocol": "..."}}],
    "security_layers": ["..."],
    "cicd_pipeline": ["..."],
    "monitoring_tools": ["..."]
}}"""
        
        analyses['tech_architecture'] = self._call_llm(tech_arch_prompt)
        
        # 4. Database Design
        if progress_callback:
            progress_callback("Designing database schema...")
        
        db_prompt = f"""Design comprehensive database schema and data models:
1. Entity-Relationship model
2. Core tables/collections
3. Relationships and constraints
4. Indexes and optimization
5. Data security and backup strategy

Project Description:
{user_input}

strictly Respond in JSON format without any marking or explanations:
{{
    "database_type": "...",
    "entities": [{{"name": "...", "description": "...", "attributes": [{{"name": "...", "type": "...", "constraints": "..."}}], "primary_key": "...", "indexes": ["..."]}}],
    "relationships": [{{"from": "...", "to": "...", "type": "one-to-many/many-to-many", "description": "..."}}],
    "optimization_strategies": ["..."],
    "backup_strategy": "...",
    "data_security": ["..."]
}}"""
        
        analyses['database_design'] = self._call_llm(db_prompt)
        
        # 5. UI/UX Design
        if progress_callback:
            progress_callback("Planning UI/UX design...")
        
        uiux_prompt = f"""Create comprehensive UI/UX design specifications:
1. User personas
2. User journeys and workflows
3. Key screens and wireframes description
4. Design system and components
5. Accessibility and responsive design

Project Description:
{user_input}

strictly Respond in JSON format without any marking or explanations:
{{
    "user_personas": [{{"name": "...", "role": "...", "goals": ["..."], "pain_points": ["..."]}}],
    "user_journeys": [{{"persona": "...", "journey": "...", "steps": ["..."], "touchpoints": ["..."]}}],
    "key_screens": [{{"screen_name": "...", "purpose": "...", "components": ["..."], "interactions": ["..."]}}],
    "design_system": {{"colors": ["..."], "typography": ["..."], "components": ["..."]}},
    "accessibility_features": ["..."],
    "responsive_breakpoints": ["..."]
}}"""
        
        analyses['uiux_design'] = self._call_llm(uiux_prompt)
        
        # 6. FRD Content (Enhanced)
        if progress_callback:
            progress_callback("Generating detailed FRD...")
        
        frd_prompt = f"""Create comprehensive Functional Requirements Document:

Project Description:
{user_input}

Please ensure all the aspects of the usecase is covered. Give me a coverage score comparing the usecase out of 100 in percentage.

Provide in JSON:
{{
    "functional_requirements": [
        {{"id": "FR-001", "requirement": "...", "priority": "High/Medium/Low", "category": "...", "description": "detailed description", "acceptance_criteria": ["..."]}}
    ],
    "use_cases": [
        {{"id": "UC-001", "title": "...", "actor": "...", "description": "...", "preconditions": ["..."], "steps": ["..."], "postconditions": ["..."], "alternate_flows": ["..."]}}
    ],
    "business_rules": [{{"rule_id": "BR-001", "rule": "...", "rationale": "..."}}],
    "data_requirements": [{{"entity": "...", "requirements": ["..."]}}],
    "interface_requirements": [{{"interface": "...", "type": "...", "requirements": ["..."]}}],
    "non_functional_requirements": [{{"category": "Performance/Security/Usability", "requirement": "...", "metric": "..."}}]
    "confidence_score": ["Please ensure all the aspects of the usecase is covered. Give me a coverage score comparing the usecase out of 100 in percentage."]

}}"""
        
        analyses['frd'] = self._call_llm(frd_prompt)
        
        # 7. Test Scripts
        if progress_callback:
            progress_callback("Creating test scripts...")
        
        test_prompt = f"""Create comprehensive functional test scripts:

Project Description:
{user_input}

Provide in JSON:
{{
    "test_strategy": "...",
    "test_scenarios": [{{"id": "TS-001", "scenario": "...", "type": "Functional/Integration/UI", "priority": "High/Medium/Low"}}],
    "test_cases": [{{"id": "TC-001", "scenario_id": "TS-001", "title": "...", "preconditions": ["..."], "steps": ["..."], "expected_results": ["..."], "test_data": "..."}}],
    "automation_candidates": ["..."],
    "performance_tests": [{{"test": "...", "criteria": "...", "expected_result": "..."}}],
    "security_tests": [{{"test": "...", "description": "..."}}]
}}"""
        
        analyses['test_scripts'] = self._call_llm(test_prompt)
        
        # 8. User Manual
        if progress_callback:
            progress_callback("Writing user manual...")
        
        manual_prompt = f"""Create comprehensive user manual content:

Project Description:
{user_input}

Provide in JSON:
{{
    "introduction": "...",
    "getting_started": [{{"step": "...", "description": "...", "screenshot_note": "..."}}],
    "features": [{{"feature": "...", "description": "...", "how_to_use": ["..."], "tips": ["..."]}}],
    "common_tasks": [{{"task": "...", "steps": ["..."], "notes": ["..."]}}],
    "troubleshooting": [{{"issue": "...", "solution": "...", "prevention": "..."}}],
    "faq": [{{"question": "...", "answer": "..."}}],
    "support_info": {{"contact": "...", "hours": "...", "resources": ["..."]}}
}}"""
        
        analyses['user_manual'] = self._call_llm(manual_prompt)
        
        return analyses
    
    def analyze_codebase(self, file_contents: Dict[str, str], progress_callback=None) -> Dict:
        """Analyze codebase using Azure OpenAI"""
        
        file_summary = "\n".join([f"File: {path}\nLines: {len(content.splitlines())}" 
                                  for path, content in list(file_contents.items())[:50]])
        
        limited_contents = {}
        total_chars = 0
        max_chars = 50000
        
        for path, content in file_contents.items():
            if total_chars + len(content) > max_chars:
                break
            limited_contents[path] = content[:5000]
            total_chars += len(content)
        
        analyses = {}
        
        # Similar structure but adapted for codebase analysis
        # Using the same prompts but with codebase context
        
        if progress_callback:
            progress_callback("Analyzing business value...")
        business_prompt = """Analyze this codebase from a BUSINESS PERSPECTIVE and provide:
1. Main business problem it solves
2. Key value propositions
3. Target users/stakeholders
4. Business benefits and ROI potential
5. Competitive advantages

Codebase structure:
{file_summary}

Sample files:
{sample_files}

strictly Respond in JSON format without any marking or explanations:
{{
    "business_problem": "...",
    "value_propositions": ["..."],
    "target_users": ["..."],
    "business_benefits": ["..."],
    "competitive_advantages": ["..."],
    "executive_summary": "..."
}}""".format(file_summary=file_summary, sample_files=self._format_files_for_prompt(limited_contents, 3))
        analyses['business'] = self._call_llm(business_prompt)
        
        if progress_callback:
            progress_callback("Analyzing application architecture...")
        app_arch_prompt = """Analyze application architecture from codebase:
1. Architecture layers and tiers
2. Application components and services
3. Communication patterns
4. Deployment architecture
5. Scalability and performance design

Files:
{sample_files}

strictly Respond in JSON format without any marking or explanations:
{{
    "architecture_layers": [{{"name": "...", "description": "...", "components": ["..."]}}],
    "services": [{{"name": "...", "type": "...", "responsibility": "...", "technologies": ["..."]}}],
    "communication_patterns": ["..."],
    "deployment_model": "...",
    "scalability_approach": "...",
    "performance_optimization": ["..."]
}}""".format(sample_files=self._format_files_for_prompt(limited_contents, 5))
        analyses['app_architecture'] = self._call_llm(app_arch_prompt)
        
        if progress_callback:
            progress_callback("Analyzing technical architecture...")
        tech_arch_prompt = """Analyze technical architecture from codebase:
1. Technology stack (frontend, backend, database, infrastructure)
2. Integration architecture
3. Security architecture
4. DevOps and CI/CD pipeline
5. Monitoring and logging

Files:
{sample_files}

strictly Respond in JSON format without any marking or explanations:
{{
    "frontend_stack": [{{"technology": "...", "purpose": "..."}}],
    "backend_stack": [{{"technology": "...", "purpose": "..."}}],
    "database_stack": [{{"technology": "...", "purpose": "..."}}],
    "infrastructure": [{{"component": "...", "technology": "...", "purpose": "..."}}],
    "integration_points": [{{"system": "...", "method": "...", "protocol": "..."}}],
    "security_layers": ["..."],
    "cicd_pipeline": ["..."],
    "monitoring_tools": ["..."]
}}""".format(sample_files=self._format_files_for_prompt(limited_contents, 5))
        analyses['tech_architecture'] = self._call_llm(tech_arch_prompt)
        
        if progress_callback:
            progress_callback("Analyzing database design...")
        db_prompt = f"""Analyze database schema from codebase:
Files: {self._format_files_for_prompt(limited_contents, 5)}
Provide in JSON: {{"database_type": "...", "entities": [{{"name": "...", "description": "...", "attributes": [{{"name": "...", "type": "...", "constraints": "..."}}], "primary_key": "...", "indexes": ["..."]}}], "relationships": [{{"from": "...", "to": "...", "type": "one-to-many", "description": "..."}}], "optimization_strategies": ["..."], "backup_strategy": "...", "data_security": ["..."]}}"""
        analyses['database_design'] = self._call_llm(db_prompt)
        
        if progress_callback:
            progress_callback("Analyzing UI/UX design...")
        uiux_prompt = f"""Analyze UI/UX design from codebase:
Files: {self._format_files_for_prompt(limited_contents, 5)}
Provide in JSON: {{"user_personas": [{{"name": "...", "role": "...", "goals": ["..."], "pain_points": ["..."]}}], "user_journeys": [{{"persona": "...", "journey": "...", "steps": ["..."], "touchpoints": ["..."]}}], "key_screens": [{{"screen_name": "...", "purpose": "...", "components": ["..."], "interactions": ["..."]}}], "design_system": {{"colors": ["..."], "typography": ["..."], "components": ["..."]}}, "accessibility_features": ["..."], "responsive_breakpoints": ["..."]}}"""
        analyses['uiux_design'] = self._call_llm(uiux_prompt)
        
        if progress_callback:
            progress_callback("Generating FRD...")
        frd_prompt = f"""Create FRD from codebase analysis:
Files: {self._format_files_for_prompt(limited_contents, 3)}
Provide in JSON: {{"functional_requirements": [{{"id": "FR-001", "requirement": "...", "priority": "High", "category": "...", "description": "...", "acceptance_criteria": ["..."]}}], "use_cases": [{{"id": "UC-001", "title": "...", "actor": "...", "description": "...", "preconditions": ["..."], "steps": ["..."], "postconditions": ["..."], "alternate_flows": ["..."]}}], "business_rules": [{{"rule_id": "BR-001", "rule": "...", "rationale": "..."}}], "data_requirements": [{{"entity": "...", "requirements": ["..."]}}], "interface_requirements": [{{"interface": "...", "type": "...", "requirements": ["..."]}}], "non_functional_requirements": [{{"category": "Performance", "requirement": "...", "metric": "..."}}]}}"""
        analyses['frd'] = self._call_llm(frd_prompt)
        
        if progress_callback:
            progress_callback("Creating test scripts...")
        test_prompt = f"""Create test scripts based on codebase:
Files: {self._format_files_for_prompt(limited_contents, 3)}
Provide in JSON: {{"test_strategy": "...", "test_scenarios": [{{"id": "TS-001", "scenario": "...", "type": "Functional", "priority": "High"}}], "test_cases": [{{"id": "TC-001", "scenario_id": "TS-001", "title": "...", "preconditions": ["..."], "steps": ["..."], "expected_results": ["..."], "test_data": "..."}}], "automation_candidates": ["..."], "performance_tests": [{{"test": "...", "criteria": "...", "expected_result": "..."}}], "security_tests": [{{"test": "...", "description": "..."}}]}}"""
        analyses['test_scripts'] = self._call_llm(test_prompt)
        
        if progress_callback:
            progress_callback("Writing user manual...")
        manual_prompt = f"""Create user manual based on codebase:
Files: {self._format_files_for_prompt(limited_contents, 3)}
Provide in JSON: {{"introduction": "...", "getting_started": [{{"step": "...", "description": "...", "screenshot_note": "..."}}], "features": [{{"feature": "...", "description": "...", "how_to_use": ["..."], "tips": ["..."]}}], "common_tasks": [{{"task": "...", "steps": ["..."], "notes": ["..."]}}], "troubleshooting": [{{"issue": "...", "solution": "...", "prevention": "..."}}], "faq": [{{"question": "...", "answer": "..."}}], "support_info": {{"contact": "...", "hours": "...", "resources": ["..."]}}}}"""
        analyses['user_manual'] = self._call_llm(manual_prompt)
        
        return analyses
    
    def _format_files_for_prompt(self, files: Dict[str, str], limit: int) -> str:
        """Format files for LLM prompt"""
        result = []
        for i, (path, content) in enumerate(list(files.items())[:limit]):
            result.append(f"\n--- {path} ---\n{content[:2000]}\n")
        return "\n".join(result)
    
    def _call_llm(self, prompt: str, max_retries: int = 3) -> Dict:
        """Call Azure OpenAI and parse JSON response with robust error handling"""
        
        for attempt in range(max_retries):
            try:
                response = self.client.chat.completions.create(
                    model=self.deployment,
                    messages=[
                        {"role": "system", "content": """You are a business and technical analyst. 
                        You MUST strictly respond with ONLY valid JSON. Do not include any markdown formatting, 
                        code blocks, or explanatory text."""},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.7,
                    max_tokens=4000,
                    response_format={"type": "json_object"}  # Force JSON mode if available
                )
                
                content = response.choices[0].message.content.strip()
                print("LLM Raw Response:", content)  # Debugging output
                # Remove markdown code blocks if present
                # if content.startswith('```'):
                #     # Extract content between ```json and ```
                #     json_match = re.search(r'```(?:json)?\s*(.*?)\s*```', content, re.DOTALL)
                #     if json_match:
                #         content = json_match.group(1).strip()
                #     else:
                #         # Remove just the ``` markers
                #         content = content.replace('```json', '').replace('```', '').strip()
                
                # # Try to find JSON object in the content
                # if not content.startswith('{'):
                #     # Try to extract JSON from the content
                #     json_match = re.search(r'\{.*\}', content, re.DOTALL)
                #     if json_match:
                #         content = json_match.group(0)
                
                # Parse JSON
                parsed_json = json.loads(content)
                return parsed_json
                
            except json.JSONDecodeError as e:
                st.warning(f"JSON parsing error (attempt {attempt + 1}/{max_retries}): {str(e)[:100]}")
                
                if attempt == max_retries - 1:
                    # Last attempt - try to fix common JSON issues
                    try:
                        # Fix unescaped quotes and newlines
                        content = content.replace('\n', ' ').replace('\r', '')
                        # Try to extract valid JSON using regex
                        json_match = re.search(r'\{[^{}]*(?:\{[^{}]*\}[^{}]*)*\}', content, re.DOTALL)
                        if json_match:
                            content = json_match.group(0)
                            parsed_json = json.loads(content)
                            return parsed_json
                    except:
                        pass
                    
                    # Return a minimal valid structure
                    st.error(f"Failed to parse JSON after {max_retries} attempts. Using fallback structure.")
                    return self._get_fallback_structure(prompt)
            
            except Exception as e:
                st.error(f"LLM call failed (attempt {attempt + 1}/{max_retries}): {str(e)}")
                if attempt == max_retries - 1:
                    return self._get_fallback_structure(prompt)
        
        return self._get_fallback_structure(prompt)


    def _get_fallback_structure(self, prompt: str) -> Dict:
        """Return a minimal fallback structure based on the prompt type"""
        
        # Determine what type of structure to return based on prompt keywords
        if 'business' in prompt.lower() or 'executive' in prompt.lower():
            return {
                "business_problem": "Analysis in progress - please review generated documents",
                "value_propositions": ["Comprehensive solution", "Scalable architecture", "User-focused design"],
                "target_users": ["End users", "Administrators", "Stakeholders"],
                "business_benefits": ["Improved efficiency", "Cost reduction", "Enhanced user experience"],
                "competitive_advantages": ["Modern technology stack", "Scalable design"],
                "executive_summary": "Comprehensive business solution with modern architecture and user-centric design."
            }
        
        elif 'application architecture' in prompt.lower() or 'architecture_layers' in prompt.lower():
            return {
                "architecture_layers": [
                    {"name": "Presentation Layer", "description": "User interface and client-side logic", "components": ["Web UI", "Mobile UI"]},
                    {"name": "Business Logic Layer", "description": "Core business rules and processing", "components": ["Services", "Controllers"]},
                    {"name": "Data Access Layer", "description": "Database operations and data management", "components": ["Repositories", "Data Models"]},
                    {"name": "Infrastructure Layer", "description": "Cross-cutting concerns and utilities", "components": ["Logging", "Security", "Caching"]}
                ],
                "services": [
                    {"name": "API Service", "type": "REST API", "responsibility": "Handle client requests", "technologies": ["REST", "HTTP"]},
                    {"name": "Business Service", "type": "Core Logic", "responsibility": "Process business rules", "technologies": ["Application Framework"]},
                    {"name": "Data Service", "type": "Data Access", "responsibility": "Manage data operations", "technologies": ["ORM", "Database"]}
                ],
                "communication_patterns": ["REST API", "Event-driven", "Request-Response"],
                "deployment_model": "Cloud-native containerized deployment",
                "scalability_approach": "Horizontal scaling with load balancing",
                "performance_optimization": ["Caching", "Connection pooling", "Async processing"]
            }
        
        elif 'technical architecture' in prompt.lower() or 'technology stack' in prompt.lower():
            return {
                "frontend_stack": [
                    {"technology": "React/Vue/Angular", "purpose": "Modern web framework"},
                    {"technology": "HTML5/CSS3", "purpose": "Markup and styling"}
                ],
                "backend_stack": [
                    {"technology": "Node.js/Python/Java", "purpose": "Server-side logic"},
                    {"technology": "REST API", "purpose": "API layer"}
                ],
                "database_stack": [
                    {"technology": "PostgreSQL/MySQL", "purpose": "Relational data storage"},
                    {"technology": "Redis", "purpose": "Caching"}
                ],
                "infrastructure": [
                    {"component": "Cloud Platform", "technology": "AWS/Azure/GCP", "purpose": "Hosting infrastructure"},
                    {"component": "Container", "technology": "Docker", "purpose": "Application containerization"}
                ],
                "integration_points": [
                    {"system": "External API", "method": "REST", "protocol": "HTTPS"}
                ],
                "security_layers": ["Authentication", "Authorization", "Encryption", "Input Validation"],
                "cicd_pipeline": ["Source Control", "Build", "Test", "Deploy"],
                "monitoring_tools": ["Application logs", "Performance metrics", "Error tracking"]
            }
        
        elif 'database' in prompt.lower():
            return {
                "database_type": "Relational Database (PostgreSQL/MySQL)",
                "entities": [
                    {
                        "name": "User",
                        "description": "User account information",
                        "attributes": [
                            {"name": "id", "type": "integer", "constraints": "PRIMARY KEY"},
                            {"name": "username", "type": "varchar(100)", "constraints": "UNIQUE NOT NULL"},
                            {"name": "email", "type": "varchar(255)", "constraints": "UNIQUE NOT NULL"},
                            {"name": "created_at", "type": "timestamp", "constraints": "DEFAULT CURRENT_TIMESTAMP"}
                        ],
                        "primary_key": "id",
                        "indexes": ["idx_username", "idx_email"]
                    }
                ],
                "relationships": [
                    {"from": "User", "to": "Profile", "type": "one-to-one", "description": "User has one profile"}
                ],
                "optimization_strategies": ["Indexing on frequently queried columns", "Query optimization", "Connection pooling"],
                "backup_strategy": "Daily automated backups with point-in-time recovery",
                "data_security": ["Encryption at rest", "Encryption in transit", "Access control"]
            }
        
        elif 'ui' in prompt.lower() or 'ux' in prompt.lower():
            return {
                "user_personas": [
                    {"name": "Primary User", "role": "End User", "goals": ["Accomplish tasks efficiently"], "pain_points": ["Complex interfaces"]}
                ],
                "user_journeys": [
                    {"persona": "Primary User", "journey": "Main workflow", "steps": ["Login", "Navigate", "Complete task"], "touchpoints": ["Web interface"]}
                ],
                "key_screens": [
                    {"screen_name": "Dashboard", "purpose": "Overview of key information", "components": ["Navigation", "Content area"], "interactions": ["Click", "Scroll"]}
                ],
                "design_system": {
                    "colors": ["Primary color", "Secondary color", "Accent color"],
                    "typography": ["Heading font", "Body font"],
                    "components": ["Buttons", "Forms", "Cards"]
                },
                "accessibility_features": ["Keyboard navigation", "Screen reader support", "Color contrast"],
                "responsive_breakpoints": ["Mobile: 320px-767px", "Tablet: 768px-1023px", "Desktop: 1024px+"]
            }
        
        elif 'frd' in prompt.lower() or 'functional requirements' in prompt.lower():
            return {
                "functional_requirements": [
                    {"id": "FR-001", "requirement": "User authentication", "priority": "High", "category": "Security", 
                    "description": "System shall provide secure user authentication", "acceptance_criteria": ["Users can login", "Passwords are encrypted"]}
                ],
                "use_cases": [
                    {"id": "UC-001", "title": "User Login", "actor": "User", "description": "User logs into the system",
                    "preconditions": ["User has account"], "steps": ["Enter credentials", "Submit form", "System validates"],
                    "postconditions": ["User is authenticated"], "alternate_flows": ["Password reset flow"]}
                ],
                "business_rules": [
                    {"rule_id": "BR-001", "rule": "Business rule description", "rationale": "Business justification"}
                ],
                "data_requirements": [
                    {"entity": "User", "requirements": ["Store user credentials", "Track login history"]}
                ],
                "interface_requirements": [
                    {"interface": "Web Interface", "type": "Browser-based", "requirements": ["Responsive design", "Modern browsers"]}
                ],
                "non_functional_requirements": [
                    {"category": "Performance", "requirement": "Page load time", "metric": "< 2 seconds"},
                    {"category": "Security", "requirement": "Data encryption", "metric": "256-bit AES"}
                ]
            }
        
        elif 'test' in prompt.lower():
            return {
                "test_strategy": "Comprehensive testing approach including functional, integration, and system testing",
                "test_scenarios": [
                    {"id": "TS-001", "scenario": "User login functionality", "type": "Functional", "priority": "High"}
                ],
                "test_cases": [
                    {"id": "TC-001", "scenario_id": "TS-001", "title": "Successful login",
                    "preconditions": ["User account exists"], "steps": ["Enter credentials", "Click login"],
                    "expected_results": ["User is logged in"], "test_data": "Valid credentials"}
                ],
                "automation_candidates": ["Login tests", "API tests"],
                "performance_tests": [
                    {"test": "Load test", "criteria": "1000 concurrent users", "expected_result": "Response time < 2s"}
                ],
                "security_tests": [
                    {"test": "Authentication test", "description": "Verify secure authentication"}
                ]
            }
        
        elif 'user manual' in prompt.lower():
            return {
                "introduction": "This user manual provides comprehensive guidance on using the system.",
                "getting_started": [
                    {"step": "Installation", "description": "Install the application", "screenshot_note": "Installation wizard"}
                ],
                "features": [
                    {"feature": "Main Feature", "description": "Primary functionality",
                    "how_to_use": ["Step 1", "Step 2"], "tips": ["Use shortcuts"]}
                ],
                "common_tasks": [
                    {"task": "Common task", "steps": ["Step 1", "Step 2"], "notes": ["Important note"]}
                ],
                "troubleshooting": [
                    {"issue": "Common issue", "solution": "Solution steps", "prevention": "How to prevent"}
                ],
                "faq": [
                    {"question": "Common question", "answer": "Detailed answer"}
                ],
                "support_info": {
                    "contact": "support@company.com",
                    "hours": "24/7",
                    "resources": ["Documentation", "Knowledge base"]
                }
            }
        
        # Default fallback
        return {
            "message": "Analysis in progress",
            "status": "Please review generated documents for details"
        }



class WordDocumentGenerator:
    """Generates professional Word documents"""
    
    def __init__(self, company_name: str = "Your Company"):
        self.company_name = company_name
    
    def _add_title_page(self, doc, title: str, project_name: str):
        """Add title page to document"""
        title_para = doc.add_heading(title, 0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_paragraph(f'\n{project_name}\n', style='Title').alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph(f'{self.company_name}', style='Subtitle').alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph(f'{datetime.now().strftime("%B %d, %Y")}', style='Subtitle').alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_page_break()
    
    def generate_frd(self, frd_data: Dict, project_name: str) -> io.BytesIO:
        """Generate Enhanced Functional Requirements Document"""
        doc = Document()
        self._add_title_page(doc, 'Functional Requirements Document', project_name)
        
        # Functional Requirements
        doc.add_heading('1. Functional Requirements', 1)
        table = doc.add_table(rows=1, cols=5)
        table.style = 'Light Grid Accent 1'
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'ID'
        hdr_cells[1].text = 'Requirement'
        hdr_cells[2].text = 'Priority'
        hdr_cells[3].text = 'Category'
        hdr_cells[4].text = 'Description'
        
        for req in frd_data.get('functional_requirements', []):
            row_cells = table.add_row().cells
            row_cells[0].text = req.get('id', '')
            row_cells[1].text = req.get('requirement', '')
            row_cells[2].text = req.get('priority', '')
            row_cells[3].text = req.get('category', '')
            row_cells[4].text = req.get('description', '')
        
        # Use Cases
        doc.add_heading('2. Use Cases', 1)
        for uc in frd_data.get('use_cases', []):
            doc.add_heading(f"{uc.get('id', '')}: {uc.get('title', '')}", 2)
            doc.add_paragraph(f"Actor: {uc.get('actor', '')}")
            doc.add_paragraph(f"Description: {uc.get('description', '')}")
            
            doc.add_heading('Preconditions:', 3)
            for pre in uc.get('preconditions', []):
                doc.add_paragraph(pre, style='List Bullet')
            
            doc.add_heading('Steps:', 3)
            for i, step in enumerate(uc.get('steps', []), 1):
                doc.add_paragraph(f'{i}. {step}', style='List Number')
            
            doc.add_heading('Postconditions:', 3)
            for post in uc.get('postconditions', []):
                doc.add_paragraph(post, style='List Bullet')
            
            if uc.get('alternate_flows'):
                doc.add_heading('Alternate Flows:', 3)
                for flow in uc.get('alternate_flows', []):
                    doc.add_paragraph(flow, style='List Bullet')
        
        # Business Rules
        doc.add_heading('3. Business Rules', 1)
        for rule in frd_data.get('business_rules', []):
            doc.add_paragraph(f"{rule.get('rule_id', '')}: {rule.get('rule', '')}")
            if rule.get('rationale'):
                doc.add_paragraph(f"Rationale: {rule.get('rationale', '')}", style='List Bullet')
        
        # Non-Functional Requirements
        doc.add_heading('4. Non-Functional Requirements', 1)
        for nfr in frd_data.get('non_functional_requirements', []):
            doc.add_paragraph(f"{nfr.get('category', '')}: {nfr.get('requirement', '')} - {nfr.get('metric', '')}", style='List Bullet')
        
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer
    
    def generate_database_design(self, db_data: Dict, project_name: str) -> io.BytesIO:
        """Generate Database Design Document"""
        doc = Document()
        self._add_title_page(doc, 'Database Design Document', project_name)
        
        # Database Type
        doc.add_heading('1. Database Overview', 1)
        doc.add_paragraph(f"Database Type: {db_data.get('database_type', 'N/A')}")
        
        # Entities
        doc.add_heading('2. Data Entities', 1)
        for entity in db_data.get('entities', []):
            doc.add_heading(entity.get('name', 'Entity'), 2)
            doc.add_paragraph(entity.get('description', ''))
            
            # Attributes table
            doc.add_heading('Attributes:', 3)
            table = doc.add_table(rows=1, cols=3)
            table.style = 'Light Grid Accent 1'
            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = 'Attribute'
            hdr_cells[1].text = 'Type'
            hdr_cells[2].text = 'Constraints'
            
            for attr in entity.get('attributes', []):
                row_cells = table.add_row().cells
                row_cells[0].text = attr.get('name', '')
                row_cells[1].text = attr.get('type', '')
                row_cells[2].text = attr.get('constraints', '')
            
            doc.add_paragraph(f"Primary Key: {entity.get('primary_key', 'N/A')}")
            
            if entity.get('indexes'):
                doc.add_paragraph("Indexes:")
                for idx in entity.get('indexes', []):
                    doc.add_paragraph(idx, style='List Bullet')
        
        # Relationships
        doc.add_heading('3. Entity Relationships', 1)
        for rel in db_data.get('relationships', []):
            doc.add_paragraph(f"{rel.get('from', '')} → {rel.get('to', '')} ({rel.get('type', '')}): {rel.get('description', '')}", style='List Bullet')
        
        # Optimization
        doc.add_heading('4. Optimization Strategies', 1)
        for opt in db_data.get('optimization_strategies', []):
            doc.add_paragraph(opt, style='List Bullet')
        
        # Security & Backup
        doc.add_heading('5. Data Security', 1)
        for sec in db_data.get('data_security', []):
            doc.add_paragraph(sec, style='List Bullet')
        
        doc.add_heading('6. Backup Strategy', 1)
        doc.add_paragraph(db_data.get('backup_strategy', 'N/A'))
        
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer
    
    def generate_uiux_design(self, uiux_data: Dict, project_name: str) -> io.BytesIO:
        """Generate UI/UX Design Document"""
        doc = Document()
        self._add_title_page(doc, 'UI/UX Design Document', project_name)
        
        # User Personas
        doc.add_heading('1. User Personas', 1)
        for persona in uiux_data.get('user_personas', []):
            doc.add_heading(persona.get('name', 'Persona'), 2)
            doc.add_paragraph(f"Role: {persona.get('role', '')}")
            doc.add_paragraph("Goals:")
            for goal in persona.get('goals', []):
                doc.add_paragraph(goal, style='List Bullet')
            doc.add_paragraph("Pain Points:")
            for pain in persona.get('pain_points', []):
                doc.add_paragraph(pain, style='List Bullet')
        
        # User Journeys
        doc.add_heading('2. User Journeys', 1)
        for journey in uiux_data.get('user_journeys', []):
            doc.add_heading(f"{journey.get('persona', '')} - {journey.get('journey', '')}", 2)
            doc.add_paragraph("Steps:")
            for step in journey.get('steps', []):
                doc.add_paragraph(step, style='List Number')
        
        # Key Screens
        doc.add_heading('3. Key Screens', 1)
        for screen in uiux_data.get('key_screens', []):
            doc.add_heading(screen.get('screen_name', 'Screen'), 2)
            doc.add_paragraph(f"Purpose: {screen.get('purpose', '')}")
            doc.add_paragraph("Components:")
            for comp in screen.get('components', []):
                doc.add_paragraph(comp, style='List Bullet')
            doc.add_paragraph("Interactions:")
            for inter in screen.get('interactions', []):
                doc.add_paragraph(inter, style='List Bullet')
        
        # Design System
        doc.add_heading('4. Design System', 1)
        design_sys = uiux_data.get('design_system', {})
        doc.add_heading('Colors:', 2)
        for color in design_sys.get('colors', []):
            doc.add_paragraph(color, style='List Bullet')
        
        doc.add_heading('Typography:', 2)
        for typo in design_sys.get('typography', []):
            doc.add_paragraph(typo, style='List Bullet')
        
        doc.add_heading('Components:', 2)
        for comp in design_sys.get('components', []):
            doc.add_paragraph(comp, style='List Bullet')
        
        # Accessibility
        doc.add_heading('5. Accessibility Features', 1)
        for feature in uiux_data.get('accessibility_features', []):
            doc.add_paragraph(feature, style='List Bullet')
        
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer
    
    def generate_test_scripts(self, test_data: Dict, project_name: str) -> io.BytesIO:
        """Generate Test Scripts Document"""
        doc = Document()
        self._add_title_page(doc, 'Functional Test Scripts', project_name)
        
        # Test Strategy
        doc.add_heading('1. Test Strategy', 1)
        doc.add_paragraph(test_data.get('test_strategy', 'N/A'))
        
        # Test Scenarios
        doc.add_heading('2. Test Scenarios', 1)
        table = doc.add_table(rows=1, cols=4)
        table.style = 'Light Grid Accent 1'
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'ID'
        hdr_cells[1].text = 'Scenario'
        hdr_cells[2].text = 'Type'
        hdr_cells[3].text = 'Priority'
        
        for scenario in test_data.get('test_scenarios', []):
            row_cells = table.add_row().cells
            row_cells[0].text = scenario.get('id', '')
            row_cells[1].text = scenario.get('scenario', '')
            row_cells[2].text = scenario.get('type', '')
            row_cells[3].text = scenario.get('priority', '')
        
        # Test Cases
        doc.add_heading('3. Test Cases', 1)
        for tc in test_data.get('test_cases', []):
            doc.add_heading(f"{tc.get('id', '')}: {tc.get('title', '')}", 2)
            doc.add_paragraph(f"Scenario: {tc.get('scenario_id', '')}")
            
            doc.add_heading('Preconditions:', 3)
            for pre in tc.get('preconditions', []):
                doc.add_paragraph(pre, style='List Bullet')
            
            doc.add_heading('Test Steps:', 3)
            for i, step in enumerate(tc.get('steps', []), 1):
                doc.add_paragraph(f'{i}. {step}', style='List Number')
            
            doc.add_heading('Expected Results:', 3)
            for result in tc.get('expected_results', []):
                doc.add_paragraph(result, style='List Bullet')
            
            if tc.get('test_data'):
                doc.add_paragraph(f"Test Data: {tc.get('test_data', '')}")
        
        # Performance Tests
        doc.add_heading('4. Performance Tests', 1)
        for perf in test_data.get('performance_tests', []):
            doc.add_paragraph(f"{perf.get('test', '')}: {perf.get('criteria', '')} → {perf.get('expected_result', '')}", style='List Bullet')
        
        # Security Tests
        doc.add_heading('5. Security Tests', 1)
        for sec in test_data.get('security_tests', []):
            doc.add_paragraph(f"{sec.get('test', '')}: {sec.get('description', '')}", style='List Bullet')
        
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer
    
    def generate_user_manual(self, manual_data: Dict, project_name: str) -> io.BytesIO:
        """Generate User Manual Document"""
        doc = Document()
        self._add_title_page(doc, 'User Manual', project_name)
        
        # Introduction
        doc.add_heading('1. Introduction', 1)
        doc.add_paragraph(manual_data.get('introduction', 'N/A'))
        
        # Getting Started
        doc.add_heading('2. Getting Started', 1)
        for i, step in enumerate(manual_data.get('getting_started', []), 1):
            doc.add_heading(f"Step {i}: {step.get('step', '')}", 2)
            doc.add_paragraph(step.get('description', ''))
            if step.get('screenshot_note'):
                doc.add_paragraph(f"[Screenshot: {step.get('screenshot_note', '')}]", style='Intense Quote')
        
        # Features
        doc.add_heading('3. Features', 1)
        for feature in manual_data.get('features', []):
            doc.add_heading(feature.get('feature', 'Feature'), 2)
            doc.add_paragraph(feature.get('description', ''))
            
            doc.add_heading('How to Use:', 3)
            for step in feature.get('how_to_use', []):
                doc.add_paragraph(step, style='List Number')
            
            if feature.get('tips'):
                doc.add_heading('Tips:', 3)
                for tip in feature.get('tips', []):
                    doc.add_paragraph(f"💡 {tip}", style='List Bullet')
        
        # Common Tasks
        doc.add_heading('4. Common Tasks', 1)
        for task in manual_data.get('common_tasks', []):
            doc.add_heading(task.get('task', 'Task'), 2)
            for step in task.get('steps', []):
                doc.add_paragraph(step, style='List Number')
            if task.get('notes'):
                doc.add_paragraph("Notes:")
                for note in task.get('notes', []):
                    doc.add_paragraph(note, style='List Bullet')
        
        # Troubleshooting
        doc.add_heading('5. Troubleshooting', 1)
        for trouble in manual_data.get('troubleshooting', []):
            doc.add_heading(f"Issue: {trouble.get('issue', '')}", 2)
            doc.add_paragraph(f"Solution: {trouble.get('solution', '')}")
            if trouble.get('prevention'):
                doc.add_paragraph(f"Prevention: {trouble.get('prevention', '')}")
        
        # FAQ
        doc.add_heading('6. Frequently Asked Questions', 1)
        for faq in manual_data.get('faq', []):
            doc.add_heading(f"Q: {faq.get('question', '')}", 2)
            doc.add_paragraph(f"A: {faq.get('answer', '')}")
        
        # Support Info
        doc.add_heading('7. Support Information', 1)
        support = manual_data.get('support_info', {})
        doc.add_paragraph(f"Contact: {support.get('contact', 'N/A')}")
        doc.add_paragraph(f"Hours: {support.get('hours', 'N/A')}")
        doc.add_paragraph("Resources:")
        for resource in support.get('resources', []):
            doc.add_paragraph(resource, style='List Bullet')
        
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer


class BusinessPPTGenerator:
    """Generates stunning business proposal presentations with 15+ slides"""
    
    # Professional color palette
    COLORS = {
        'primary': RGBColor(20, 33, 61),
        'secondary': RGBColor(0, 123, 255),
        'accent': RGBColor(255, 107, 53),
        'success': RGBColor(40, 167, 69),
        'gold': RGBColor(255, 193, 7),
        'purple': RGBColor(111, 66, 193),
        'teal': RGBColor(23, 162, 184),
        'pink': RGBColor(232, 62, 140),
        'light': RGBColor(248, 249, 250),
        'text': RGBColor(33, 37, 41),
        'text_light': RGBColor(108, 117, 125),
        'border': RGBColor(222, 226, 230),
        'white': RGBColor(255, 255, 255),
    }
    
    def __init__(self, company_name: str = "Your Company"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.company_name = company_name
        
    def _add_premium_background(self, slide, style='default'):
        """Add premium background with modern design"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['white']
        
        if style == 'cover':
            # Diagonal gradient effect with shapes
            shape1 = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(7), Inches(7.5)
            )
            fill = shape1.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['primary']
            shape1.line.fill.background()
            
            shape2 = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(7), Inches(0), Inches(6.333), Inches(7.5)
            )
            fill = shape2.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['secondary']
            shape2.line.fill.background()
            
            # Decorative circles
            for x, y, size, color in [
                (9.5, -0.5, 3, self.COLORS['accent']),
                (11, 5.5, 2.5, self.COLORS['gold'])
            ]:
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(x), Inches(y), Inches(size), Inches(size)
                )
                fill = circle.fill
                fill.solid()
                fill.fore_color.rgb = color
                circle.fill.transparency = 0.2
                circle.line.fill.background()
            
        elif style == 'section':
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(13.333), Inches(7.5)
            )
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['primary']
            bg.line.fill.background()
            
            # Decorative elements
            for i, (x, y, size, color) in enumerate([
                (10, -1, 4, self.COLORS['secondary']),
                (-1, 5, 3.5, self.COLORS['accent']),
                (11, 6, 2.5, self.COLORS['gold'])
            ]):
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(x), Inches(y), Inches(size), Inches(size)
                )
                fill = circle.fill
                fill.solid()
                fill.fore_color.rgb = color
                circle.fill.transparency = 0.3
                circle.line.fill.background()
        
        else:  # content slides
            # Accent bar
            accent_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(0.15), Inches(7.5)
            )
            fill = accent_bar.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['secondary']
            accent_bar.line.fill.background()
            
            # Top right decoration
            top_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(11.5), Inches(-0.5), Inches(2.5), Inches(2)
            )
            fill = top_shape.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['light']
            top_shape.line.fill.background()
    
    def add_cover_slide(self, title: str, subtitle: str):
        """Add stunning cover slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide, 'cover')
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(5.5), Inches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        title_para.line_spacing = 1.1
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(5.5), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(22)
        subtitle_para.font.color.rgb = self.COLORS['white']
        subtitle_para.font.italic = True
        
        # Company name
        company_box = slide.shapes.add_textbox(Inches(7.5), Inches(3), Inches(5), Inches(0.6))
        company_frame = company_box.text_frame
        company_frame.text = self.company_name
        company_para = company_frame.paragraphs[0]
        company_para.font.size = Pt(28)
        company_para.font.bold = True
        company_para.font.color.rgb = self.COLORS['white']
        
        # Date
        date_box = slide.shapes.add_textbox(Inches(7.5), Inches(3.7), Inches(5), Inches(0.4))
        date_frame = date_box.text_frame
        date_frame.text = datetime.now().strftime("%B %Y")
        date_para = date_frame.paragraphs[0]
        date_para.font.size = Pt(18)
        date_para.font.color.rgb = self.COLORS['white']
    
    def add_section_divider(self, title: str, subtitle: str, icon: str = ""):
        """Add beautiful section divider"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide, 'section')
        
        # Icon
        if icon:
            icon_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(2.5), Inches(1.5))
            icon_frame = icon_box.text_frame
            icon_frame.text = icon
            icon_para = icon_frame.paragraphs[0]
            icon_para.font.size = Pt(100)
            icon_para.alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(2), Inches(3.8), Inches(9.333), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(3), Inches(5), Inches(7.333), Inches(0.6))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(20)
            subtitle_para.font.color.rgb = self.COLORS['white']
            subtitle_para.alignment = PP_ALIGN.CENTER
    
    def add_executive_summary(self, summary: str, highlights: List[str]):
        """Add executive summary slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide)
        
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(0.5), Inches(12), Inches(0.7)
        )
        fill = header.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['primary']
        header.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.6), Inches(11), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "📋 Executive Summary"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        
        # Summary text
        summary_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(1.8))
        summary_frame = summary_box.text_frame
        summary_frame.word_wrap = True
        summary_frame.text = summary
        summary_para = summary_frame.paragraphs[0]
        summary_para.font.size = Pt(18)
        summary_para.font.color.rgb = self.COLORS['text']
        summary_para.line_spacing = 1.4
        
        # Key highlights
        highlight_label = slide.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(12), Inches(0.4))
        hl_frame = highlight_label.text_frame
        hl_frame.text = "Key Highlights"
        hl_para = hl_frame.paragraphs[0]
        hl_para.font.size = Pt(22)
        hl_para.font.bold = True
        hl_para.font.color.rgb = self.COLORS['primary']
        
        colors = [self.COLORS['secondary'], self.COLORS['success'], self.COLORS['accent'], self.COLORS['purple']]
        icons = ['🎯', '💡', '🚀', '⭐']
        
        for i, highlight in enumerate(highlights[:4]):
            col = i % 2
            row = i // 2
            left = 0.7 + (col * 6.2)
            top = 4.2 + (row * 1.4)
            
            # Card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top), Inches(5.8), Inches(1.2)
            )
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['white']
            card.line.color.rgb = colors[i]
            card.line.width = Pt(2)
            
            # Icon
            icon_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left + 0.15), Inches(top + 0.25), Inches(0.7), Inches(0.7)
            )
            fill = icon_circle.fill
            fill.solid()
            fill.fore_color.rgb = colors[i]
            icon_circle.line.fill.background()
            
            icon_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.25), Inches(0.7), Inches(0.7))
            icon_frame = icon_box.text_frame
            icon_frame.text = icons[i]
            icon_para = icon_frame.paragraphs[0]
            icon_para.font.size = Pt(28)
            icon_para.alignment = PP_ALIGN.CENTER
            
            # Text
            text_box = slide.shapes.add_textbox(Inches(left + 1), Inches(top + 0.15), Inches(4.6), Inches(0.9))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame.text = highlight
            text_para = text_frame.paragraphs[0]
            text_para.font.size = Pt(14)
            text_para.font.color.rgb = self.COLORS['text']
    
    def add_frd_overview_slide(self, frd_data: Dict):
        """Add FRD Overview slide with requirements summary"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "📑 Functional Requirements Overview"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['primary']
        
        # Requirements summary stats
        reqs = frd_data.get('functional_requirements', [])
        high_priority = len([r for r in reqs if r.get('priority') == 'High'])
        med_priority = len([r for r in reqs if r.get('priority') == 'Medium'])
        low_priority = len([r for r in reqs if r.get('priority') == 'Low'])
        
        stats = [
            ('Total Requirements', len(reqs), self.COLORS['secondary']),
            ('High Priority', high_priority, self.COLORS['accent']),
            ('Medium Priority', med_priority, self.COLORS['gold']),
            ('Low Priority', low_priority, self.COLORS['success'])
        ]
        
        for i, (label, value, color) in enumerate(stats):
            left = 0.7 + (i * 3.1)
            
            # Stat card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(1.5), Inches(2.9), Inches(1.5)
            )
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = color
            fill.transparency = 0.1
            card.line.color.rgb = color
            card.line.width = Pt(3)
            
            # Value
            val_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(1.7), Inches(2.5), Inches(0.7))
            val_frame = val_box.text_frame
            val_frame.text = str(value)
            val_para = val_frame.paragraphs[0]
            val_para.font.size = Pt(48)
            val_para.font.bold = True
            val_para.font.color.rgb = color
            val_para.alignment = PP_ALIGN.CENTER
            
            # Label
            lbl_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(2.5), Inches(2.5), Inches(0.4))
            lbl_frame = lbl_box.text_frame
            lbl_frame.text = label
            lbl_para = lbl_frame.paragraphs[0]
            lbl_para.font.size = Pt(14)
            lbl_para.font.color.rgb = self.COLORS['text']
            lbl_para.alignment = PP_ALIGN.CENTER
        
        # Top requirements list
        req_label = slide.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(12), Inches(0.4))
        req_frame = req_label.text_frame
        req_frame.text = "Key Functional Requirements"
        req_para = req_frame.paragraphs[0]
        req_para.font.size = Pt(20)
        req_para.font.bold = True
        req_para.font.color.rgb = self.COLORS['primary']
        
        # List top 6 requirements
        for i, req in enumerate(reqs[:6]):
            row = i // 2
            col = i % 2
            left = 0.7 + (col * 6.2)
            top = 3.9 + (row * 0.9)
            
            priority_colors = {'High': self.COLORS['accent'], 'Medium': self.COLORS['gold'], 'Low': self.COLORS['success']}
            color = priority_colors.get(req.get('priority', 'Medium'), self.COLORS['text_light'])
            
            # Requirement card
            req_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top), Inches(5.8), Inches(0.75)
            )
            fill = req_card.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['white']
            req_card.line.color.rgb = color
            req_card.line.width = Pt(2)
            
            # Priority badge
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left + 0.1), Inches(top + 0.15), Inches(0.8), Inches(0.45)
            )
            fill = badge.fill
            fill.solid()
            fill.fore_color.rgb = color
            badge.line.fill.background()
            
            badge_text = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.15), Inches(0.8), Inches(0.45))
            badge_frame = badge_text.text_frame
            badge_frame.text = req.get('priority', 'M')[0]
            badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            badge_para = badge_frame.paragraphs[0]
            badge_para.font.size = Pt(16)
            badge_para.font.bold = True
            badge_para.font.color.rgb = self.COLORS['white']
            badge_para.alignment = PP_ALIGN.CENTER
            
            # Requirement text
            text_box = slide.shapes.add_textbox(Inches(left + 1.1), Inches(top + 0.1), Inches(4.5), Inches(0.55))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame.text = f"{req.get('id', '')}: {req.get('requirement', '')[:60]}..."
            text_para = text_frame.paragraphs[0]
            text_para.font.size = Pt(11)
            text_para.font.color.rgb = self.COLORS['text']
    
    def add_application_architecture_slide(self, app_arch_data: Dict):
        """Add Application Architecture slide with layers and services"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "🏗️ Application Architecture"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['primary']
        
        # Architecture layers
        layers = app_arch_data.get('architecture_layers', [])[:4]
        layer_colors = [self.COLORS['accent'], self.COLORS['secondary'], self.COLORS['success'], self.COLORS['purple']]
        
        for i, layer in enumerate(layers):
            top = 1.5 + (i * 1.4)
            
            # Layer box
            layer_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.7), Inches(top), Inches(5.8), Inches(1.2)
            )
            fill = layer_box.fill
            fill.solid()
            fill.fore_color.rgb = layer_colors[i]
            fill.transparency = 0.15
            layer_box.line.color.rgb = layer_colors[i]
            layer_box.line.width = Pt(3)
            
            # Layer name
            name_box = slide.shapes.add_textbox(Inches(0.9), Inches(top + 0.1), Inches(5.4), Inches(0.4))
            name_frame = name_box.text_frame
            name_frame.text = layer.get('name', 'Layer')
            name_para = name_frame.paragraphs[0]
            name_para.font.size = Pt(18)
            name_para.font.bold = True
            name_para.font.color.rgb = layer_colors[i]
            
            # Description
            desc_box = slide.shapes.add_textbox(Inches(0.9), Inches(top + 0.5), Inches(5.4), Inches(0.6))
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            desc_frame.text = layer.get('description', '')
            desc_para = desc_frame.paragraphs[0]
            desc_para.font.size = Pt(12)
            desc_para.font.color.rgb = self.COLORS['text']
        
        # Key Services section
        services_label = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.6), Inches(0.4))
        serv_frame = services_label.text_frame
        serv_frame.text = "Key Services & Components"
        serv_para = serv_frame.paragraphs[0]
        serv_para.font.size = Pt(20)
        serv_para.font.bold = True
        serv_para.font.color.rgb = self.COLORS['primary']
        
        # Services list
        services = app_arch_data.get('services', [])[:5]
        for i, service in enumerate(services):
            top = 2.1 + (i * 0.95)
            
            # Service card
            serv_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(7), Inches(top), Inches(5.6), Inches(0.85)
            )
            fill = serv_card.fill
            fill.solid()
            fill.fore_color.rgb = self.COLORS['light']
            serv_card.line.color.rgb = self.COLORS['border']
            serv_card.line.width = Pt(1)
            
            # Service name
            name_box = slide.shapes.add_textbox(Inches(7.2), Inches(top + 0.1), Inches(5.2), Inches(0.3))
            name_frame = name_box.text_frame
            name_frame.text = f"🔹 {service.get('name', 'Service')}"
            name_para = name_frame.paragraphs[0]
            name_para.font.size = Pt(14)
            name_para.font.bold = True
            name_para.font.color.rgb = self.COLORS['secondary']
            
            # Service description
            desc_box = slide.shapes.add_textbox(Inches(7.2), Inches(top + 0.4), Inches(5.2), Inches(0.4))
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            desc_frame.text = service.get('responsibility', '')[:80]
            desc_para = desc_frame.paragraphs[0]
            desc_para.font.size = Pt(10)
            desc_para.font.color.rgb = self.COLORS['text_light']
    
    def add_technical_architecture_slide(self, tech_arch_data: Dict):
        """Add Technical Architecture slide with tech stack"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "⚙️ Technical Architecture"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['primary']
        
        # Tech stack sections
        sections = [
            ('Frontend Stack', tech_arch_data.get('frontend_stack', [])[:4], self.COLORS['accent'], 0.7, 1.5),
            ('Backend Stack', tech_arch_data.get('backend_stack', [])[:4], self.COLORS['secondary'], 7, 1.5),
            ('Database Stack', tech_arch_data.get('database_stack', [])[:4], self.COLORS['success'], 0.7, 4.2),
            ('Infrastructure', tech_arch_data.get('infrastructure', [])[:4], self.COLORS['purple'], 7, 4.2)
        ]
        
        for section_name, items, color, left, top in sections:
            # Section header
            header = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top), Inches(5.6), Inches(0.5)
            )
            fill = header.fill
            fill.solid()
            fill.fore_color.rgb = color
            header.line.fill.background()
            
            header_text = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.05), Inches(5.2), Inches(0.4))
            header_frame = header_text.text_frame
            header_frame.text = section_name
            header_para = header_frame.paragraphs[0]
            header_para.font.size = Pt(16)
            header_para.font.bold = True
            header_para.font.color.rgb = self.COLORS['white']
            
            # Items
            for i, item in enumerate(items):
                item_top = top + 0.6 + (i * 0.45)
                
                item_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(item_top), Inches(5.2), Inches(0.4))
                item_frame = item_box.text_frame
                item_frame.word_wrap = True
                
                tech_name = item.get('technology', item.get('component', 'Technology'))
                purpose = item.get('purpose', '')[:50]
                item_frame.text = f"▪ {tech_name}"
                
                item_para = item_frame.paragraphs[0]
                item_para.font.size = Pt(12)
                item_para.font.color.rgb = self.COLORS['text']
                
                if purpose:
                    purpose_para = item_frame.add_paragraph()
                    purpose_para.text = f"  {purpose}"
                    purpose_para.font.size = Pt(9)
                    purpose_para.font.color.rgb = self.COLORS['text_light']
    
    def add_database_design_slide(self, db_data: Dict):
        """Add Database Design slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "🗄️ Database Design & Data Models"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['primary']
        
        # Database type
        db_type_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(1.4), Inches(5.8), Inches(0.6)
        )
        fill = db_type_box.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['secondary']
        db_type_box.line.fill.background()
        
        db_type_text = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5.2), Inches(0.4))
        db_type_frame = db_type_text.text_frame
        db_type_frame.text = f"Database Type: {db_data.get('database_type', 'N/A')}"
        db_type_para = db_type_frame.paragraphs[0]
        db_type_para.font.size = Pt(18)
        db_type_para.font.bold = True
        db_type_para.font.color.rgb = self.COLORS['white']
        
        # Core entities
        entities_label = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(5.8), Inches(0.3))
        ent_frame = entities_label.text_frame
        ent_frame.text = "Core Entities"
        ent_para = ent_frame.paragraphs[0]
        ent_para.font.size = Pt(18)
        ent_para.font.bold = True
        ent_para.font.color.rgb = self.COLORS['primary']
        
        entities = db_data.get('entities', [])[:6]
        entity_colors = [self.COLORS['accent'], self.COLORS['teal'], self.COLORS['gold'], 
                        self.COLORS['pink'], self.COLORS['success'], self.COLORS['purple']]
        
        for i, entity in enumerate(entities):
            row = i // 2
            col = i % 2
            left = 0.7 + (col * 3)
            top = 2.6 + (row * 1.3)
            
            # Entity card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top), Inches(2.7), Inches(1.1)
            )
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = entity_colors[i % len(entity_colors)]
            fill.transparency = 0.15
            card.line.color.rgb = entity_colors[i % len(entity_colors)]
            card.line.width = Pt(2)
            
            # Entity name
            name_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.1), Inches(2.4), Inches(0.35))
            name_frame = name_box.text_frame
            name_frame.text = entity.get('name', 'Entity')
            name_para = name_frame.paragraphs[0]
            name_para.font.size = Pt(14)
            name_para.font.bold = True
            name_para.font.color.rgb = entity_colors[i % len(entity_colors)]
            
            # Attributes count
            attrs = entity.get('attributes', [])
            attr_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.5), Inches(2.4), Inches(0.5))
            attr_frame = attr_box.text_frame
            attr_frame.word_wrap = True
            attr_frame.text = f"📊 {len(attrs)} attributes\n🔑 PK: {entity.get('primary_key', 'N/A')}"
            attr_para = attr_frame.paragraphs[0]
            attr_para.font.size = Pt(10)
            attr_para.font.color.rgb = self.COLORS['text']
        
        # Relationships & Optimization
        rel_label = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.6), Inches(0.3))
        rel_frame = rel_label.text_frame
        rel_frame.text = "Entity Relationships"
        rel_para = rel_frame.paragraphs[0]
        rel_para.font.size = Pt(18)
        rel_para.font.bold = True
        rel_para.font.color.rgb = self.COLORS['primary']
        
        relationships = db_data.get('relationships', [])[:5]
        for i, rel in enumerate(relationships):
            top = 1.9 + (i * 0.55)
            
            rel_box = slide.shapes.add_textbox(Inches(7), Inches(top), Inches(5.6), Inches(0.5))
            rel_frame = rel_box.text_frame
            rel_frame.word_wrap = True
            rel_frame.text = f"➜ {rel.get('from', '')} → {rel.get('to', '')} ({rel.get('type', '')})"
            rel_para = rel_frame.paragraphs[0]
            rel_para.font.size = Pt(11)
            rel_para.font.color.rgb = self.COLORS['text']
        
        # Optimization strategies
        opt_label = slide.shapes.add_textbox(Inches(7), Inches(4.8), Inches(5.6), Inches(0.3))
        opt_frame = opt_label.text_frame
        opt_frame.text = "Optimization & Security"
        opt_para = opt_frame.paragraphs[0]
        opt_para.font.size = Pt(16)
        opt_para.font.bold = True
        opt_para.font.color.rgb = self.COLORS['success']
        
        strategies = db_data.get('optimization_strategies', [])[:3] + db_data.get('data_security', [])[:2]
        for i, strategy in enumerate(strategies):
            top = 5.2 + (i * 0.4)
            
            strat_box = slide.shapes.add_textbox(Inches(7), Inches(top), Inches(5.6), Inches(0.35))
            strat_frame = strat_box.text_frame
            strat_frame.word_wrap = True
            strat_frame.text = f"✓ {strategy[:80]}"
            strat_para = strat_frame.paragraphs[0]
            strat_para.font.size = Pt(10)
            strat_para.font.color.rgb = self.COLORS['text']
    
    def add_uiux_design_slide(self, uiux_data: Dict):
        """Add UI/UX Design slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "🎨 UI/UX Design & User Experience"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['primary']
        
        # User Personas
        personas_label = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(5.8), Inches(0.3))
        pers_frame = personas_label.text_frame
        pers_frame.text = "👥 User Personas"
        pers_para = pers_frame.paragraphs[0]
        pers_para.font.size = Pt(18)
        pers_para.font.bold = True
        pers_para.font.color.rgb = self.COLORS['primary']
        
        personas = uiux_data.get('user_personas', [])[:3]
        persona_colors = [self.COLORS['accent'], self.COLORS['teal'], self.COLORS['purple']]
        
        for i, persona in enumerate(personas):
            top = 1.85 + (i * 1.5)
            
            # Persona card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.7), Inches(top), Inches(5.8), Inches(1.3)
            )
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = persona_colors[i]
            fill.transparency = 0.1
            card.line.color.rgb = persona_colors[i]
            card.line.width = Pt(2)
            
            # Persona name
            name_box = slide.shapes.add_textbox(Inches(0.9), Inches(top + 0.1), Inches(5.4), Inches(0.35))
            name_frame = name_box.text_frame
            name_frame.text = f"👤 {persona.get('name', 'User')} - {persona.get('role', '')}"
            name_para = name_frame.paragraphs[0]
            name_para.font.size = Pt(14)
            name_para.font.bold = True
            name_para.font.color.rgb = persona_colors[i]
            
            # Goals
            goals_box = slide.shapes.add_textbox(Inches(0.9), Inches(top + 0.5), Inches(5.4), Inches(0.7))
            goals_frame = goals_box.text_frame
            goals_frame.word_wrap = True
            goals = persona.get('goals', [])[:2]
            goals_frame.text = "Goals: " + ", ".join(goals)
            goals_para = goals_frame.paragraphs[0]
            goals_para.font.size = Pt(11)
            goals_para.font.color.rgb = self.COLORS['text']
        
        # Key Screens
        screens_label = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.6), Inches(0.3))
        scr_frame = screens_label.text_frame
        scr_frame.text = "🖥️ Key Screens & Interfaces"
        scr_para = scr_frame.paragraphs[0]
        scr_para.font.size = Pt(18)
        scr_para.font.bold = True
        scr_para.font.color.rgb = self.COLORS['primary']
        
        screens = uiux_data.get('key_screens', [])[:6]
        screen_colors = [self.COLORS['secondary'], self.COLORS['accent'], self.COLORS['success'], 
                        self.COLORS['gold'], self.COLORS['teal'], self.COLORS['pink']]
        
        for i, screen in enumerate(screens):
            row = i // 2
            col = i % 2
            left = 7 + (col * 2.9)
            top = 1.85 + (row * 1.5)
            
            # Screen card
            screen_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top), Inches(2.7), Inches(1.3)
            )
            fill = screen_card.fill
            fill.solid()
            fill.fore_color.rgb = screen_colors[i % len(screen_colors)]
            fill.transparency = 0.15
            screen_card.line.color.rgb = screen_colors[i % len(screen_colors)]
            screen_card.line.width = Pt(2)
            
            # Screen name
            scr_name_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.1), Inches(2.4), Inches(0.35))
            scr_name_frame = scr_name_box.text_frame
            scr_name_frame.word_wrap = True
            scr_name_frame.text = screen.get('screen_name', 'Screen')
            scr_name_para = scr_name_frame.paragraphs[0]
            scr_name_para.font.size = Pt(13)
            scr_name_para.font.bold = True
            scr_name_para.font.color.rgb = screen_colors[i % len(screen_colors)]
            
            # Purpose
            purpose_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.5), Inches(2.4), Inches(0.7))
            purpose_frame = purpose_box.text_frame
            purpose_frame.word_wrap = True
            purpose_frame.text = screen.get('purpose', '')[:80]
            purpose_para = purpose_frame.paragraphs[0]
            purpose_para.font.size = Pt(10)
            purpose_para.font.color.rgb = self.COLORS['text']
    
    def add_test_strategy_slide(self, test_data: Dict):
        """Add Testing Strategy slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "🧪 Testing Strategy & Test Scripts"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['primary']
        
        # Test Strategy
        strategy_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(1))
        strategy_frame = strategy_box.text_frame
        strategy_frame.word_wrap = True
        strategy_frame.text = test_data.get('test_strategy', 'Comprehensive testing approach covering functional, integration, performance, and security testing.')
        strategy_para = strategy_frame.paragraphs[0]
        strategy_para.font.size = Pt(16)
        strategy_para.font.color.rgb = self.COLORS['text']
        strategy_para.line_spacing = 1.3
        
        # Test Scenarios Summary
        scenarios = test_data.get('test_scenarios', [])
        func_tests = len([s for s in scenarios if s.get('type') == 'Functional'])
        int_tests = len([s for s in scenarios if s.get('type') == 'Integration'])
        ui_tests = len([s for s in scenarios if s.get('type') == 'UI'])
        
        stats = [
            ('Total Scenarios', len(scenarios), self.COLORS['secondary']),
            ('Functional Tests', func_tests, self.COLORS['accent']),
            ('Integration Tests', int_tests, self.COLORS['success']),
            ('UI Tests', ui_tests, self.COLORS['purple'])
        ]
        
        for i, (label, value, color) in enumerate(stats):
            left = 0.7 + (i * 3.1)
            
            # Stat card
            stat_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(2.6), Inches(2.9), Inches(1.2)
            )
            fill = stat_card.fill
            fill.solid()
            fill.fore_color.rgb = color
            fill.transparency = 0.1
            stat_card.line.color.rgb = color
            stat_card.line.width = Pt(3)
            
            # Value
            val_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(2.7), Inches(2.5), Inches(0.6))
            val_frame = val_box.text_frame
            val_frame.text = str(value)
            val_para = val_frame.paragraphs[0]
            val_para.font.size = Pt(42)
            val_para.font.bold = True
            val_para.font.color.rgb = color
            val_para.alignment = PP_ALIGN.CENTER
            
            # Label
            lbl_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(3.3), Inches(2.5), Inches(0.4))
            lbl_frame = lbl_box.text_frame
            lbl_frame.text = label
            lbl_para = lbl_frame.paragraphs[0]
            lbl_para.font.size = Pt(12)
            lbl_para.font.color.rgb = self.COLORS['text']
            lbl_para.alignment = PP_ALIGN.CENTER
        
        # Key Test Areas
        areas_label = slide.shapes.add_textbox(Inches(0.7), Inches(4.1), Inches(12), Inches(0.3))
        areas_frame = areas_label.text_frame
        areas_frame.text = "Key Testing Areas"
        areas_para = areas_frame.paragraphs[0]
        areas_para.font.size = Pt(18)
        areas_para.font.bold = True
        areas_para.font.color.rgb = self.COLORS['primary']
        
        # Performance & Security tests
        perf_tests = test_data.get('performance_tests', [])[:3]
        sec_tests = test_data.get('security_tests', [])[:3]
        
        col1_tests = perf_tests
        col2_tests = sec_tests
        
        # Column 1: Performance
        perf_label = slide.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(5.8), Inches(0.3))
        perf_frame = perf_label.text_frame
        perf_frame.text = "⚡ Performance Tests"
        perf_para = perf_frame.paragraphs[0]
        perf_para.font.size = Pt(14)
        perf_para.font.bold = True
        perf_para.font.color.rgb = self.COLORS['gold']
        
        for i, test in enumerate(col1_tests):
            top = 5 + (i * 0.6)
            test_box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(5.8), Inches(0.55))
            test_frame = test_box.text_frame
            test_frame.word_wrap = True
            test_frame.text = f"▪ {test.get('test', '')}: {test.get('criteria', '')}"
            test_para = test_frame.paragraphs[0]
            test_para.font.size = Pt(11)
            test_para.font.color.rgb = self.COLORS['text']
        
        # Column 2: Security
        sec_label = slide.shapes.add_textbox(Inches(7), Inches(4.6), Inches(5.6), Inches(0.3))
        sec_frame = sec_label.text_frame
        sec_frame.text = "🔒 Security Tests"
        sec_para = sec_frame.paragraphs[0]
        sec_para.font.size = Pt(14)
        sec_para.font.bold = True
        sec_para.font.color.rgb = self.COLORS['accent']
        
        for i, test in enumerate(col2_tests):
            top = 5 + (i * 0.6)
            test_box = slide.shapes.add_textbox(Inches(7), Inches(top), Inches(5.6), Inches(0.55))
            test_frame = test_box.text_frame
            test_frame.word_wrap = True
            test_frame.text = f"▪ {test.get('test', '')}"
            test_para = test_frame.paragraphs[0]
            test_para.font.size = Pt(11)
            test_para.font.color.rgb = self.COLORS['text']
    
    def add_user_manual_overview_slide(self, manual_data: Dict):
        """Add User Manual Overview slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_premium_background(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "📖 User Manual & Documentation"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['primary']
        
        # Introduction
        intro_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(1))
        intro_frame = intro_box.text_frame
        intro_frame.word_wrap = True
        intro_frame.text = manual_data.get('introduction', 'Comprehensive user guide covering all features and functionality.')
        intro_para = intro_frame.paragraphs[0]
        intro_para.font.size = Pt(16)
        intro_para.font.color.rgb = self.COLORS['text']
        intro_para.line_spacing = 1.3
        
        # Key Features
        features_label = slide.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(12), Inches(0.3))
        feat_frame = features_label.text_frame
        feat_frame.text = "Key Features Covered"
        feat_para = feat_frame.paragraphs[0]
        feat_para.font.size = Pt(20)
        feat_para.font.bold = True
        feat_para.font.color.rgb = self.COLORS['primary']
        
        features = manual_data.get('features', [])[:6]
        feature_colors = [self.COLORS['secondary'], self.COLORS['accent'], self.COLORS['success'],
                         self.COLORS['purple'], self.COLORS['teal'], self.COLORS['gold']]
        
        for i, feature in enumerate(features):
            row = i // 2
            col = i % 2
            left = 0.7 + (col * 6.2)
            top = 3.1 + (row * 1.1)
            
            # Feature card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top), Inches(5.8), Inches(0.95)
            )
            fill = card.fill
            fill.solid()
            fill.fore_color.rgb = feature_colors[i]
            fill.transparency = 0.1
            card.line.color.rgb = feature_colors[i]
            card.line.width = Pt(2)
            
            # Feature name
            name_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.1), Inches(5.5), Inches(0.35))
            name_frame = name_box.text_frame
            name_frame.text = f"📌 {feature.get('feature', 'Feature')}"
            name_para = name_frame.paragraphs[0]
            name_para.font.size = Pt(14)
            name_para.font.bold = True
            name_para.font.color.rgb = feature_colors[i]
            
            # Description
            desc_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.5), Inches(5.5), Inches(0.4))
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            desc_frame.text = feature.get('description', '')[:80]
            desc_para = desc_frame.paragraphs[0]
            desc_para.font.size = Pt(11)
            desc_para.font.color.rgb = self.COLORS['text']
        
        # Support section
        support_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(6.4), Inches(12), Inches(0.8)
        )
        fill = support_box.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['light']
        support_box.line.color.rgb = self.COLORS['border']
        support_box.line.width = Pt(2)
        
        support_text = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.5), Inches(0.6))
        support_frame = support_text.text_frame
        support_info = manual_data.get('support_info', {})
        support_frame.text = f"📞 Support: {support_info.get('contact', 'support@company.com')} | {support_info.get('hours', '24/7')}"
        support_para = support_frame.paragraphs[0]
        support_para.font.size = Pt(14)
        support_para.font.color.rgb = self.COLORS['text']
        support_para.alignment = PP_ALIGN.CENTER
    
    def save(self, filename: str = "proposal.pptx") -> io.BytesIO:
        """Save presentation to BytesIO"""
        buffer = io.BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer


# Main Streamlit UI
def main():
    st.title("🚀 Business Proposal & Architecture Generator")
    st.markdown("Generate comprehensive business proposals, technical documentation, and architecture diagrams from your codebase or project description")
    
    
    # Sidebar configuration
    with st.sidebar:
        st.header("⚙️ Configuration")
        
        # Azure OpenAI settings
        azure_endpoint = st.text_input("Azure OpenAI Endpoint", type="password")
        azure_key = st.text_input("Azure OpenAI Key", type="password")
        azure_deployment = st.text_input("Deployment Name", value="gpt-4")
        company_name = st.text_input("Company Name", value="Your Company")
        project_name = st.text_input("Project Name", value="My Project")
        
        st.divider()
        
        # Input mode selection
        st.subheader("📥 Input Method")
        input_mode = st.radio(
            "Choose input method:",
            ["Upload Codebase (ZIP)", "Describe Project (Text)"],
            key="input_mode_radio"
        )
        
        st.session_state.input_mode = 'upload' if input_mode == "Upload Codebase (ZIP)" else 'text'
    
    # Main content area
    if not azure_endpoint or not azure_key:
        st.warning("⚠️ Please configure Azure OpenAI credentials in the sidebar")
        return
    
    # Input section based on mode
    if st.session_state.input_mode == 'upload':
        st.header("📦 Upload Your Codebase")
        st.markdown("Upload a ZIP file containing your project code")
        
        uploaded_file = st.file_uploader("Choose a ZIP file", type=['zip'])
        
        if uploaded_file and st.button("🔍 Analyze Codebase", type="primary"):
            analyze_codebase_flow(uploaded_file, azure_endpoint, azure_key, azure_deployment, company_name, project_name)
    
    else:  # text input mode
        st.header("✏️ Describe Your Project")
        st.markdown("Provide a detailed description of your project, including features, goals, and technical requirements")
        
        user_input = st.text_area(
            "Project Description",
            height=300,
            placeholder="Example: We want to build an e-commerce platform that allows users to browse products, add items to cart, make secure payments, and track orders. The system should support multiple vendors, real-time inventory management, and provide analytics dashboards for business insights..."
        )
        
        if user_input and st.button("🔍 Analyze Requirements", type="primary"):
            analyze_text_flow(user_input, azure_endpoint, azure_key, azure_deployment, company_name, project_name)
    
    # Display results and download options
    if st.session_state.analysis_complete and st.session_state.analysis_data:
        st.success("✅ Analysis Complete!")
        
        display_analysis_results(st.session_state.analysis_data, azure_endpoint, azure_key, azure_deployment)

        
        st.divider()
        st.header("📥 Download Documents")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.subheader("📊 Presentation")
            if st.button("Generate PowerPoint (15+ Slides)", type="primary", use_container_width=True):
                with st.spinner("Creating stunning presentation..."):
                    ppt_buffer = generate_presentation(st.session_state.analysis_data, company_name, project_name)
                    st.download_button(
                        label="⬇️ Download PowerPoint",
                        data=ppt_buffer,
                        file_name=f"{project_name}_Proposal.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
        
        with col2:
            st.subheader("📄 Word Documents")
            
            doc_gen = WordDocumentGenerator(company_name)
            
            # FRD
            if 'frd' in st.session_state.analysis_data:
                frd_buffer = doc_gen.generate_frd(st.session_state.analysis_data['frd'], project_name)
                st.download_button(
                    label="⬇️ Download FRD",
                    data=frd_buffer,
                    file_name=f"{project_name}_FRD.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True
                )
            
            # Database Design
            if 'database_design' in st.session_state.analysis_data:
                db_buffer = doc_gen.generate_database_design(st.session_state.analysis_data['database_design'], project_name)
                st.download_button(
                    label="⬇️ Download Database Design",
                    data=db_buffer,
                    file_name=f"{project_name}_DatabaseDesign.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True
                )
            
            # UI/UX Design
            if 'uiux_design' in st.session_state.analysis_data:
                uiux_buffer = doc_gen.generate_uiux_design(st.session_state.analysis_data['uiux_design'], project_name)
                st.download_button(
                    label="⬇️ Download UI/UX Design",
                    data=uiux_buffer,
                    file_name=f"{project_name}_UIUX_Design.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True
                )
            
            # Test Scripts
            if 'test_scripts' in st.session_state.analysis_data:
                test_buffer = doc_gen.generate_test_scripts(st.session_state.analysis_data['test_scripts'], project_name)
                st.download_button(
                    label="⬇️ Download Test Scripts",
                    data=test_buffer,
                    file_name=f"{project_name}_TestScripts.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True
                )
            
            # User Manual
            if 'user_manual' in st.session_state.analysis_data:
                manual_buffer = doc_gen.generate_user_manual(st.session_state.analysis_data['user_manual'], project_name)
                st.download_button(
                    label="⬇️ Download User Manual",
                    data=manual_buffer,
                    file_name=f"{project_name}_UserManual.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True
                )

def display_traceability_tab(azure_endpoint, azure_key, azure_deployment):
    """Single tab for traceability metrics analysis"""
    
    st.subheader("🔍 Requirements Traceability Analysis")
    st.markdown("Upload your code and requirement documents to analyze alignment and generate a traceability score")
    
    # Upload section
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("##### 📁 Upload Code Files")
        code_files = st.file_uploader(
            "Code files (ZIP or individual files)",
            type=['zip', 'py', 'js', 'jsx', 'ts', 'tsx', 'java', 'cpp', 'cs', 'go', 'rb', 'php'],
            accept_multiple_files=True,
            key="trace_code",
            help="Upload your source code files or a ZIP containing your codebase"
        )
    
    with col2:
        st.markdown("##### 📄 Upload Requirement Documents")
        doc_files = st.file_uploader(
            "Requirement documents (PDF, DOCX, TXT)",
            type=['pdf', 'docx', 'txt', 'doc'],
            accept_multiple_files=True,
            key="trace_docs",
            help="Upload FRD, BRD, or other requirement documents"
        )
    
    # Analyze button
    if code_files and doc_files:
        if st.button("🔬 Analyze Traceability", type="primary", use_container_width=True):
            
            # Extract code content
            with st.spinner("📥 Extracting code content..."):
                try:
                    code_content = []
                    for file in code_files:
                        if file.name.endswith('.zip'):
                            with zipfile.ZipFile(file, 'r') as zip_ref:
                                for file_info in zip_ref.filelist[:20]:
                                    if not file_info.is_dir():
                                        try:
                                            content = zip_ref.read(file_info.filename).decode('utf-8', errors='ignore')
                                            code_content.append(f"=== {file_info.filename} ===\n{content[:3000]}\n")
                                        except:
                                            continue
                        else:
                            content = file.read().decode('utf-8', errors='ignore')
                            code_content.append(f"=== {file.name} ===\n{content[:3000]}\n")
                    
                    code_text = "\n".join(code_content)
                    st.success(f"✅ Extracted {len(code_text)} characters from {len(code_files)} code file(s)")
                except Exception as e:
                    st.error(f"❌ Error extracting code: {e}")
                    return
            
            # Extract document content
            with st.spinner("📥 Extracting document content..."):
                try:
                    doc_content = []
                    for file in doc_files:
                        if file.name.endswith('.txt'):
                            content = file.read().decode('utf-8', errors='ignore')
                            doc_content.append(f"=== {file.name} ===\n{content[:5000]}\n")
                        
                        elif file.name.endswith('.docx'):
                            doc = Document(file)
                            text = "\n".join([para.text for para in doc.paragraphs])
                            doc_content.append(f"=== {file.name} ===\n{text[:5000]}\n")
                        
                        elif file.name.endswith('.pdf'):
                            try:
                                import PyPDF2
                                pdf_reader = PyPDF2.PdfReader(file)
                                text = ""
                                for page in pdf_reader.pages[:10]:
                                    text += page.extract_text()
                                doc_content.append(f"=== {file.name} ===\n{text[:5000]}\n")
                            except ImportError:
                                st.warning(f"⚠️ PyPDF2 not installed. Skipping {file.name}")
                            except Exception as e:
                                st.warning(f"⚠️ Could not read PDF {file.name}: {e}")
                    
                    doc_text = "\n".join(doc_content)
                    st.success(f"✅ Extracted {len(doc_text)} characters from {len(doc_files)} document(s)")
                except Exception as e:
                    st.error(f"❌ Error extracting documents: {e}")
                    return
            
            # Analyze traceability
            with st.spinner("🔬 Analyzing traceability with AI... This may take a moment..."):
                try:
                    # Call LLM for analysis
                    prompt = f"""You are a requirements traceability analyst. Analyze the alignment between the provided code and requirement documents.

CODE IMPLEMENTATION:
{code_text[:15000]}

REQUIREMENT DOCUMENTS:
{doc_text[:15000]}

Provide a comprehensive traceability analysis in JSON format:
{{
    "overall_score": 85,
    "traceability_matrix": [
        {{
            "requirement_id": "FR-001",
            "requirement": "User authentication",
            "implementation_status": "Implemented",
            "code_references": ["auth.py:45-89"],
            "coverage_score": 95,
            "notes": "Fully implemented with JWT"
        }}
    ],
    "coverage_analysis": {{
        "requirements_implemented": 15,
        "requirements_partial": 3,
        "requirements_missing": 2,
        "total_requirements": 20,
        "implementation_percentage": 85
    }},
    "quality_metrics": {{
        "code_completeness": 85,
        "documentation_alignment": 78,
        "test_coverage_alignment": 65,
        "architecture_compliance": 90
    }},
    "gaps_identified": [
        {{
            "gap_type": "Missing Feature",
            "requirement": "FR-005: Password reset",
            "severity": "High",
            "recommendation": "Implement password reset flow"
        }}
    ],
    "implemented_but_not_documented": [
        {{
            "feature": "Two-factor authentication",
            "code_location": "auth.py:120-156",
            "recommendation": "Add to requirements"
        }}
    ],
    "recommendations": [
        "Implement missing password reset feature",
        "Add unit tests for authentication"
    ]
}}"""

                    client = AzureOpenAI(
                        azure_endpoint=azure_endpoint,
                        api_key=azure_key,
                        api_version="2024-05-01-preview"
                    )
                    
                    response = client.chat.completions.create(
                        model=azure_deployment,
                        messages=[
                            {"role": "system", "content": "You are a requirements traceability expert. Respond only with valid JSON."},
                            {"role": "user", "content": prompt}
                        ],
                        temperature=0.3,
                        max_tokens=4000,
                        response_format={"type": "json_object"}
                    )
                    
                    results = json.loads(response.choices[0].message.content.strip())
                    
                except Exception as e:
                    st.error(f"❌ Analysis failed: {e}")
                    return
            
            # Display Results
            st.success("✅ Analysis Complete!")
            st.divider()
            
            # === OVERALL SCORE ===
            st.header("🎯 Traceability Score")
            overall_score = results.get('overall_score', 0)
            
            score_col1, score_col2, score_col3 = st.columns([2, 1, 1])
            
            with score_col1:
                st.metric("Overall Traceability Score", f"{overall_score}%", 
                        help="Overall alignment between code and requirements")
                
                if overall_score >= 85:
                    st.success("✅ Excellent - Strong alignment!")
                elif overall_score >= 70:
                    st.info("✓ Good - Minor gaps exist")
                elif overall_score >= 50:
                    st.warning("⚠️ Moderate - Several gaps")
                else:
                    st.error("❌ Poor - Major gaps found")
            
            coverage = results.get('coverage_analysis', {})
            with score_col2:
                impl = coverage.get('requirements_implemented', 0)
                total = coverage.get('total_requirements', 0)
                st.metric("Requirements Met", f"{impl}/{total}")
            
            with score_col3:
                impl_pct = coverage.get('implementation_percentage', 0)
                st.metric("Implementation %", f"{impl_pct}%")
            
            st.divider()
            
            # === DETAILED SECTIONS ===
            detail_tabs = st.tabs([
                "📊 Coverage", 
                "📋 Matrix", 
                "⚠️ Gaps", 
                "📈 Quality",
                "💡 Actions"
            ])
            
            # TAB 1: Coverage Analysis
            with detail_tabs[0]:
                st.subheader("Requirements Coverage Breakdown")
                
                cov_col1, cov_col2, cov_col3, cov_col4 = st.columns(4)
                
                with cov_col1:
                    st.metric("✅ Implemented", coverage.get('requirements_implemented', 0))
                with cov_col2:
                    st.metric("🔶 Partial", coverage.get('requirements_partial', 0))
                with cov_col3:
                    st.metric("❌ Missing", coverage.get('requirements_missing', 0))
                with cov_col4:
                    st.metric("📝 Total", coverage.get('total_requirements', 0))
                
                # Visual progress
                st.markdown("##### Coverage Progress")
                if total > 0:
                    st.progress(impl / total)
                    st.caption(f"{impl} out of {total} requirements fully implemented")
            
            # TAB 2: Traceability Matrix
            with detail_tabs[1]:
                st.subheader("Requirements Traceability Matrix")
                
                matrix = results.get('traceability_matrix', [])
                
                if matrix:
                    for item in matrix:
                        status = item.get('implementation_status', 'Unknown')
                        
                        if status == 'Implemented':
                            icon = "✅"
                        elif status == 'Partial':
                            icon = "🔶"
                        else:
                            icon = "❌"
                        
                        with st.expander(f"{icon} {item.get('requirement_id', 'N/A')} - {item.get('requirement', 'N/A')}"):
                            mat_col1, mat_col2 = st.columns([3, 1])
                            
                            with mat_col1:
                                st.write(f"**Status:** {status}")
                                st.write(f"**Notes:** {item.get('notes', 'N/A')}")
                                
                                if item.get('code_references'):
                                    st.write("**Code References:**")
                                    for ref in item.get('code_references', []):
                                        st.code(ref, language="text")
                            
                            with mat_col2:
                                cov_score = item.get('coverage_score', 0)
                                st.metric("Coverage", f"{cov_score}%")
                                st.progress(cov_score / 100)
                else:
                    st.info("No detailed matrix available")
            
            # TAB 3: Gaps & Issues
            with detail_tabs[2]:
                st.subheader("Identified Gaps and Issues")
                
                gaps = results.get('gaps_identified', [])
                
                if gaps:
                    high = [g for g in gaps if g.get('severity') == 'High']
                    medium = [g for g in gaps if g.get('severity') == 'Medium']
                    low = [g for g in gaps if g.get('severity') == 'Low']
                    
                    if high:
                        st.error(f"🔴 High Severity ({len(high)})")
                        for gap in high:
                            with st.expander(f"❗ {gap.get('requirement', 'N/A')}"):
                                st.write(f"**Type:** {gap.get('gap_type', 'N/A')}")
                                st.write(f"**Recommendation:** {gap.get('recommendation', 'N/A')}")
                    
                    if medium:
                        st.warning(f"🟡 Medium Severity ({len(medium)})")
                        for gap in medium:
                            with st.expander(f"⚠️ {gap.get('requirement', 'N/A')}"):
                                st.write(f"**Type:** {gap.get('gap_type', 'N/A')}")
                                st.write(f"**Recommendation:** {gap.get('recommendation', 'N/A')}")
                    
                    if low:
                        st.info(f"🔵 Low Severity ({len(low)})")
                        for gap in low:
                            with st.expander(f"ℹ️ {gap.get('requirement', 'N/A')}"):
                                st.write(f"**Type:** {gap.get('gap_type', 'N/A')}")
                                st.write(f"**Recommendation:** {gap.get('recommendation', 'N/A')}")
                else:
                    st.success("✅ No significant gaps identified!")
                
                # Undocumented implementations
                st.markdown("##### 📌 Implemented but Not Documented")
                undoc = results.get('implemented_but_not_documented', [])
                
                if undoc:
                    for item in undoc:
                        with st.expander(f"💡 {item.get('feature', 'N/A')}"):
                            st.write(f"**Location:** {item.get('code_location', 'N/A')}")
                            st.write(f"**Action:** {item.get('recommendation', 'N/A')}")
                else:
                    st.info("All implementations are documented")
            
            # TAB 4: Quality Metrics
            with detail_tabs[3]:
                st.subheader("Quality Metrics Dashboard")
                
                metrics = results.get('quality_metrics', {})
                
                qual_col1, qual_col2 = st.columns(2)
                
                with qual_col1:
                    st.markdown("**Code Completeness**")
                    completeness = metrics.get('code_completeness', 0)
                    st.progress(completeness / 100)
                    st.caption(f"{completeness}%")
                    
                    st.markdown("**Documentation Alignment**")
                    doc_align = metrics.get('documentation_alignment', 0)
                    st.progress(doc_align / 100)
                    st.caption(f"{doc_align}%")
                
                with qual_col2:
                    st.markdown("**Test Coverage Alignment**")
                    test_cov = metrics.get('test_coverage_alignment', 0)
                    st.progress(test_cov / 100)
                    st.caption(f"{test_cov}%")
                    
                    st.markdown("**Architecture Compliance**")
                    arch_comp = metrics.get('architecture_compliance', 0)
                    st.progress(arch_comp / 100)
                    st.caption(f"{arch_comp}%")
                
                # Average
                avg_quality = sum(metrics.values()) / len(metrics) if metrics else 0
                
                st.divider()
                st.metric("Average Quality Score", f"{avg_quality:.1f}%")
                
                if avg_quality >= 85:
                    st.success("🌟 Excellent quality!")
                elif avg_quality >= 70:
                    st.info("👍 Good quality")
                elif avg_quality >= 50:
                    st.warning("⚠️ Needs improvement")
                else:
                    st.error("❌ Quality issues detected")
            
            # TAB 5: Recommendations
            with detail_tabs[4]:
                st.subheader("💡 Recommended Actions")
                
                recommendations = results.get('recommendations', [])
                
                if recommendations:
                    for i, rec in enumerate(recommendations, 1):
                        st.markdown(f"**{i}.** {rec}")
                else:
                    st.success("✅ No actions needed - excellent alignment!")
                
                # Download report
                st.divider()
                st.markdown("##### 📄 Export Report")
                
                report = f"""
REQUIREMENTS TRACEABILITY ANALYSIS REPORT
{'='*80}

OVERALL SCORE: {overall_score}%

COVERAGE ANALYSIS
{'-'*80}
Total Requirements: {coverage.get('total_requirements', 0)}
Implemented: {coverage.get('requirements_implemented', 0)}
Partial: {coverage.get('requirements_partial', 0)}
Missing: {coverage.get('requirements_missing', 0)}
Implementation %: {coverage.get('implementation_percentage', 0)}%

QUALITY METRICS
{'-'*80}
Code Completeness: {metrics.get('code_completeness', 0)}%
Documentation Alignment: {metrics.get('documentation_alignment', 0)}%
Test Coverage: {metrics.get('test_coverage_alignment', 0)}%
Architecture Compliance: {metrics.get('architecture_compliance', 0)}%

GAPS IDENTIFIED
{'-'*80}
"""
                
                for gap in gaps:
                    report += f"\n[{gap.get('severity')}] {gap.get('requirement')}\n"
                    report += f"  Type: {gap.get('gap_type')}\n"
                    report += f"  Action: {gap.get('recommendation')}\n"
                
                report += f"\nRECOMMENDATIONS\n{'-'*80}\n"
                for i, rec in enumerate(recommendations, 1):
                    report += f"{i}. {rec}\n"
                
                report += f"\n{'='*80}\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
                
                st.download_button(
                    label="📥 Download Full Report (TXT)",
                    data=report,
                    file_name=f"traceability_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                    mime="text/plain",
                    use_container_width=True
                )
    
    else:
        st.info("👆 Upload code files and requirement documents to begin analysis")

def analyze_codebase_flow(uploaded_file, azure_endpoint, azure_key, azure_deployment, company_name, project_name):
    """Flow for analyzing uploaded codebase"""
    with tempfile.TemporaryDirectory() as temp_dir:
        # Save uploaded file
        zip_path = os.path.join(temp_dir, "codebase.zip")
        with open(zip_path, 'wb') as f:
            f.write(uploaded_file.getvalue())
        
        # Extract and analyze
        analyzer = CodebaseAnalyzer(azure_endpoint, azure_key, azure_deployment)
        
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        def update_progress(message):
            status_text.text(message)
        
        try:
            status_text.text("Extracting files...")
            progress_bar.progress(10)
            file_contents = analyzer.extract_files(zip_path, temp_dir)
            
            st.info(f"🔍 Found {len(file_contents)} code files")
            
            progress_bar.progress(20)
            status_text.text("Analyzing codebase with AI...")
            
            analyses = analyzer.analyze_codebase(file_contents, update_progress)
            
            progress_bar.progress(100)
            status_text.text("Analysis complete!")
            
            st.session_state.analysis_data = analyses
            st.session_state.analysis_complete = True
            
        except Exception as e:
            st.error(f"❌ Error during analysis: {e}")
            import traceback
            st.code(traceback.format_exc())


def analyze_text_flow(user_input, azure_endpoint, azure_key, azure_deployment, company_name, project_name):
    """Flow for analyzing text input"""
    analyzer = CodebaseAnalyzer(azure_endpoint, azure_key, azure_deployment)
    
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    def update_progress(message):
        status_text.text(message)
        # Estimate progress
        progress_map = {
            "Analyzing business": 12,
            "Designing application architecture": 24,
            "Designing technical architecture": 36,
            "Designing database": 48,
            "Planning UI/UX": 60,
            "Generating detailed FRD": 72,
            "Creating test scripts": 84,
            "Writing user manual": 96
        }
        for key, value in progress_map.items():
            if key in message:
                progress_bar.progress(value)
                break
    
    try:
        progress_bar.progress(5)
        status_text.text("Starting analysis...")
        
        analyses = analyzer.analyze_from_text(user_input, update_progress)
        
        progress_bar.progress(100)
        status_text.text("Analysis complete!")
        
        st.session_state.analysis_data = analyses
        st.session_state.analysis_complete = True
        
    except Exception as e:
        st.error(f"❌ Error during analysis: {e}")
        import traceback
        st.code(traceback.format_exc())


def display_analysis_results(analyses, azure_endpoint, azure_key, azure_deployment):
    """Display analysis results in tabs"""
    st.header("📊 Analysis Results")
    
    tabs = st.tabs([
        "Executive Summary", 
        "FRD",
        "Application Architecture",
        "Technical Architecture",
        "Database Design",
        "UI/UX Design",
        "Test Scripts",
        "User Manual",
        "Tranceability Matrix"
    ])
    
    # Executive Summary Tab
    with tabs[0]:
        business = analyses.get('business', {})
        st.subheader("Executive Summary")
        st.write(business.get('executive_summary', 'N/A'))
        
        st.subheader("Business Problem")
        st.write(business.get('business_problem', 'N/A'))
        
        col1, col2 = st.columns(2)
        with col1:
            st.subheader("Target Users")
            for user in business.get('target_users', []):
                st.write(f"• {user}")
        
        with col2:
            st.subheader("Business Benefits")
            for benefit in business.get('business_benefits', []):
                st.write(f"• {benefit}")
        
        st.subheader("Value Propositions")
        for i, prop in enumerate(business.get('value_propositions', []), 1):
            st.markdown(f"**{i}.** {prop}")
    
    # FRD Tab
    with tabs[1]:
        frd = analyses.get('frd', {})
        
        st.subheader("Functional Requirements")
        for req in frd.get('functional_requirements', []):
            priority_color = {'High': '🔴', 'Medium': '🟡', 'Low': '🟢'}.get(req.get('priority', ''), '⚪')
            with st.expander(f"{priority_color} {req.get('id', '')} - {req.get('requirement', '')}"):
                st.write(f"**Category:** {req.get('category', 'N/A')}")
                st.write(f"**Priority:** {req.get('priority', 'N/A')}")
                st.write(f"**Description:** {req.get('description', 'N/A')}")
                if req.get('acceptance_criteria'):
                    st.write("**Acceptance Criteria:**")
                    for criterion in req.get('acceptance_criteria', []):
                        st.write(f"✓ {criterion}")
        
        st.subheader("Use Cases")
        for uc in frd.get('use_cases', []):
            with st.expander(f"{uc.get('id', '')} - {uc.get('title', '')}"):
                st.write(f"**Actor:** {uc.get('actor', 'N/A')}")
                st.write(f"**Description:** {uc.get('description', 'N/A')}")
                
                st.write("**Steps:**")
                for i, step in enumerate(uc.get('steps', []), 1):
                    st.write(f"{i}. {step}")
        
        st.subheader("Non-Functional Requirements")
        for nfr in frd.get('non_functional_requirements', []):
            st.write(f"**{nfr.get('category', '')}:** {nfr.get('requirement', '')} - {nfr.get('metric', '')}")
    
    # Application Architecture Tab
    with tabs[2]:
        app_arch = analyses.get('app_architecture', {})
        
        st.subheader("Architecture Layers")
        for layer in app_arch.get('architecture_layers', []):
            with st.expander(f"🏗️ {layer.get('name', 'Layer')}"):
                st.write(layer.get('description', 'N/A'))
                st.write("**Components:**")
                for comp in layer.get('components', []):
                    st.write(f"• {comp}")
        
        st.subheader("Services")
        for service in app_arch.get('services', []):
            st.write(f"**{service.get('name', 'Service')}** ({service.get('type', 'N/A')})")
            st.write(f"→ {service.get('responsibility', 'N/A')}")
        
        st.subheader("Communication Patterns")
        for pattern in app_arch.get('communication_patterns', []):
            st.write(f"• {pattern}")
    
    # Technical Architecture Tab
    with tabs[3]:
        tech_arch = analyses.get('tech_architecture', {})
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.subheader("Frontend Stack")
            for tech in tech_arch.get('frontend_stack', []):
                st.write(f"**{tech.get('technology', 'Tech')}:** {tech.get('purpose', '')}")
            
            st.subheader("Backend Stack")
            for tech in tech_arch.get('backend_stack', []):
                st.write(f"**{tech.get('technology', 'Tech')}:** {tech.get('purpose', '')}")
        
        with col2:
            st.subheader("Database Stack")
            for tech in tech_arch.get('database_stack', []):
                st.write(f"**{tech.get('technology', 'Tech')}:** {tech.get('purpose', '')}")
            
            st.subheader("Infrastructure")
            for infra in tech_arch.get('infrastructure', []):
                st.write(f"**{infra.get('component', 'Component')}:** {infra.get('technology', '')} - {infra.get('purpose', '')}")
        
        st.subheader("Integration Points")
        for integration in tech_arch.get('integration_points', []):
            st.write(f"• {integration.get('system', '')}: {integration.get('method', '')} via {integration.get('protocol', '')}")
    
    # Database Design Tab
    with tabs[4]:
        db_data = analyses.get('database_design', {})
        
        st.info(f"Database Type: {db_data.get('database_type', 'N/A')}")
        
        st.subheader("Data Entities")
        for entity in db_data.get('entities', []):
            with st.expander(f"📊 {entity.get('name', 'Entity')}"):
                st.write(entity.get('description', ''))
                st.write(f"**Primary Key:** {entity.get('primary_key', 'N/A')}")
                
                if entity.get('attributes'):
                    st.write("**Attributes:**")
                    for attr in entity.get('attributes', []):
                        st.write(f"• {attr.get('name', '')}: {attr.get('type', '')} - {attr.get('constraints', '')}")
        
        st.subheader("Relationships")
        for rel in db_data.get('relationships', []):
            st.write(f"➜ {rel.get('from', '')} → {rel.get('to', '')} ({rel.get('type', '')}): {rel.get('description', '')}")
        
        st.subheader("Optimization Strategies")
        for opt in db_data.get('optimization_strategies', []):
            st.write(f"✓ {opt}")
    
    # UI/UX Design Tab
    with tabs[5]:
        uiux_data = analyses.get('uiux_design', {})
        
        st.subheader("User Personas")
        for persona in uiux_data.get('user_personas', []):
            with st.expander(f"👤 {persona.get('name', 'User')} - {persona.get('role', '')}"):
                st.write("**Goals:**")
                for goal in persona.get('goals', []):
                    st.write(f"• {goal}")
                st.write("**Pain Points:**")
                for pain in persona.get('pain_points', []):
                    st.write(f"• {pain}")
        
        st.subheader("Key Screens")
        for screen in uiux_data.get('key_screens', []):
            with st.expander(f"🖥️ {screen.get('screen_name', 'Screen')}"):
                st.write(f"**Purpose:** {screen.get('purpose', '')}")
                st.write("**Components:**")
                for comp in screen.get('components', []):
                    st.write(f"• {comp}")
                st.write("**Interactions:**")
                for inter in screen.get('interactions', []):
                    st.write(f"• {inter}")
        
        st.subheader("Design System")
        design_sys = uiux_data.get('design_system', {})
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.write("**Colors:**")
            for color in design_sys.get('colors', []):
                st.write(f"• {color}")
        
        with col2:
            st.write("**Typography:**")
            for typo in design_sys.get('typography', []):
                st.write(f"• {typo}")
        
        with col3:
            st.write("**Components:**")
            for comp in design_sys.get('components', []):
                st.write(f"• {comp}")
    
    # Test Scripts Tab
    with tabs[6]:
        test_data = analyses.get('test_scripts', {})
        
        st.subheader("Test Strategy")
        st.write(test_data.get('test_strategy', 'N/A'))
        
        st.subheader("Test Scenarios")
        for scenario in test_data.get('test_scenarios', []):
            priority_color = {'High': '🔴', 'Medium': '🟡', 'Low': '🟢'}.get(scenario.get('priority', ''), '⚪')
            st.write(f"{priority_color} **{scenario.get('id', '')}:** {scenario.get('scenario', '')} ({scenario.get('type', '')})")
        
        st.subheader("Test Cases")
        for tc in test_data.get('test_cases', []):
            with st.expander(f"{tc.get('id', '')} - {tc.get('title', '')}"):
                st.write(f"**Scenario:** {tc.get('scenario_id', '')}")
                st.write("**Steps:**")
                for i, step in enumerate(tc.get('steps', []), 1):
                    st.write(f"{i}. {step}")
                st.write("**Expected Results:**")
                for result in tc.get('expected_results', []):
                    st.write(f"✓ {result}")
    
    # User Manual Tab
    with tabs[7]:
        manual_data = analyses.get('user_manual', {})
        
        st.subheader("Introduction")
        st.write(manual_data.get('introduction', 'N/A'))
        
        st.subheader("Getting Started")
        for i, step in enumerate(manual_data.get('getting_started', []), 1):
            st.write(f"**Step {i}: {step.get('step', '')}**")
            st.write(step.get('description', ''))
        
        st.subheader("Features")
        for feature in manual_data.get('features', []):
            with st.expander(f"📌 {feature.get('feature', 'Feature')}"):
                st.write(feature.get('description', ''))
                st.write("**How to Use:**")
                for step in feature.get('how_to_use', []):
                    st.write(f"• {step}")
                if feature.get('tips'):
                    st.write("**Tips:**")
                    for tip in feature.get('tips', []):
                        st.write(f"💡 {tip}")
        
        st.subheader("Support Information")
        support = manual_data.get('support_info', {})
        st.info(f"📞 Contact: {support.get('contact', 'N/A')} | Hours: {support.get('hours', 'N/A')}")

    with tabs[8]:
        display_traceability_tab( azure_endpoint, azure_key, azure_deployment)
def generate_presentation(analyses, company_name, project_name):
    """Generate comprehensive 15+ slide PowerPoint presentation"""
    ppt = BusinessPPTGenerator(company_name)
    
    # Slide 1: Cover
    ppt.add_cover_slide(
        title=project_name,
        subtitle="Comprehensive Business Proposal & Technical Architecture"
    )
    
    # Slide 2: Executive Summary Section
    ppt.add_section_divider("Executive Summary", "Business Overview & Strategic Value", "📋")
    
    # Slide 3: Executive Summary Content
    business = analyses.get('business', {})
    ppt.add_executive_summary(
        summary=business.get('executive_summary', 'N/A'),
        highlights=business.get('value_propositions', [])[:4]
    )
    
    # Slide 4: FRD Section
    ppt.add_section_divider("Functional Requirements", "Detailed Requirements Analysis", "📑")
    
    # Slide 5: FRD Overview
    frd = analyses.get('frd', {})
    ppt.add_frd_overview_slide(frd)
    
    # Slide 6: Application Architecture Section
    ppt.add_section_divider("Application Architecture", "System Design & Components", "🏗️")
    
    # Slide 7: Application Architecture Details
    app_arch = analyses.get('app_architecture', {})
    ppt.add_application_architecture_slide(app_arch)
    
    # Slide 8: Technical Architecture Section
    ppt.add_section_divider("Technical Architecture", "Technology Stack & Infrastructure", "⚙️")
    
    # Slide 9: Technical Architecture Details
    tech_arch = analyses.get('tech_architecture', {})
    ppt.add_technical_architecture_slide(tech_arch)
    
    # Slide 10: Database Design Section
    ppt.add_section_divider("Database Design", "Data Models & Schema", "🗄️")
    
    # Slide 11: Database Design Details
    db_data = analyses.get('database_design', {})
    ppt.add_database_design_slide(db_data)
    
    # Slide 12: UI/UX Design Section
    ppt.add_section_divider("UI/UX Design", "User Experience & Interface", "🎨")
    
    # Slide 13: UI/UX Design Details
    uiux_data = analyses.get('uiux_design', {})
    ppt.add_uiux_design_slide(uiux_data)
    
    # Slide 14: Testing Section
    ppt.add_section_divider("Testing Strategy", "Quality Assurance & Test Scripts", "🧪")
    
    # Slide 15: Testing Details
    test_data = analyses.get('test_scripts', {})
    ppt.add_test_strategy_slide(test_data)
    
    # Slide 16: User Manual Section
    ppt.add_section_divider("User Documentation", "Guides & Support", "📖")
    
    # Slide 17: User Manual Overview
    manual_data = analyses.get('user_manual', {})
    ppt.add_user_manual_overview_slide(manual_data)
    
    return ppt.save()


if __name__ == "__main__":
    main()
